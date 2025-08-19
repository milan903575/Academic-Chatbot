import os
import json
import numpy as np
from datetime import datetime, timedelta
from typing import List, Dict, Any, Optional, Tuple, Union
from dataclasses import dataclass, asdict
import time
import io
import base64
from pathlib import Path
import re
from collections import defaultdict, Counter
import math
import asyncio
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed
import threading
from queue import Queue, Empty
import multiprocessing as mp
from functools import partial
import signal
import sys
import shutil
import hashlib

# LangChain imports
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import (
    PyPDFLoader, UnstructuredPDFLoader, Docx2txtLoader,
    UnstructuredWordDocumentLoader, UnstructuredExcelLoader,
    UnstructuredPowerPointLoader, TextLoader, BSHTMLLoader
)
from langchain.text_splitter import RecursiveCharacterTextSplitter, TokenTextSplitter
from langchain.schema import Document
from langchain.memory import ConversationSummaryBufferMemory
from langchain.prompts import PromptTemplate, ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough

# Document processing imports
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract
from pdf2image import convert_from_path
import pandas as pd
from dotenv import load_dotenv
import tiktoken
from pdfminer.high_level import extract_text
from pdfminer.layout import LAParams
import pdfplumber
from bs4 import BeautifulSoup

# Advanced search imports
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from sklearn.cluster import KMeans
from sklearn.decomposition import LatentDirichletAllocation
from sklearn.feature_extraction.text import CountVectorizer

# Load environment variables
load_dotenv()

# Enhanced NLTK setup with comprehensive fallbacks
def setup_nltk_safely():
    """Setup NLTK with comprehensive error handling and fallbacks"""
    try:
        import nltk
        from nltk.tokenize import word_tokenize, sent_tokenize
        from nltk.corpus import stopwords
        from nltk.stem import WordNetLemmatizer

        datasets_to_download = [
            ('punkt', 'tokenizers/punkt'),
            ('punkt_tab', 'tokenizers/punkt_tab'),
            ('stopwords', 'corpora/stopwords'),
            ('wordnet', 'corpora/wordnet'),
            ('omw-1.4', 'corpora/omw-1.4')
        ]

        successful_downloads = []
        for dataset, path in datasets_to_download:
            try:
                nltk.data.find(path)
                successful_downloads.append(dataset)
            except LookupError:
                try:
                    nltk.download(dataset, quiet=True, raise_on_error=False)
                    try:
                        nltk.data.find(path)
                        successful_downloads.append(dataset)
                    except:
                        pass
                except Exception as e:
                    pass

        global STOP_WORDS, LEMMATIZER, safe_word_tokenize, safe_lemmatize

        if 'stopwords' in successful_downloads:
            try:
                STOP_WORDS = set(stopwords.words('english'))
            except Exception:
                STOP_WORDS = set(['the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should'])
        else:
            STOP_WORDS = set(['the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should'])

        if 'wordnet' in successful_downloads:
            try:
                LEMMATIZER = WordNetLemmatizer()
            except Exception:
                LEMMATIZER = None
        else:
            LEMMATIZER = None

        def safe_word_tokenize(text: str) -> List[str]:
            """Safely tokenize text with fallback"""
            if 'punkt' in successful_downloads or 'punkt_tab' in successful_downloads:
                try:
                    return word_tokenize(text)
                except Exception:
                    pass
            return re.findall(r'\b\w+\b', text.lower())

        def safe_lemmatize(word: str) -> str:
            """Safely lemmatize word with fallback"""
            if LEMMATIZER:
                try:
                    return LEMMATIZER.lemmatize(word)
                except Exception:
                    pass
            return word

        return True

    except ImportError:
        STOP_WORDS = set(['the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should'])
        LEMMATIZER = None

        def safe_word_tokenize(text: str) -> List[str]:
            return re.findall(r'\b\w+\b', text.lower())

        def safe_lemmatize(word: str) -> str:
            return word

        return False

# Initialize NLTK
setup_nltk_safely()

@dataclass
class ConversationContext:
    """Enhanced conversation context with memory"""
    session_id: str
    conversation_history: List[Dict[str, str]]
    last_query: str
    last_response: str
    context_documents: List[str]
    query_intent: str
    follow_up_context: Dict[str, Any]
    timestamp: datetime
    total_exchanges: int
    immediate_context_weight: float = 2.0

@dataclass
class GroupConversationContext:
    """Group-specific conversation context"""
    group_id: int
    group_name: str
    thread_id: int
    thread_title: str
    member_count: int
    conversation_history: List[Dict[str, str]]
    last_query: str
    last_response: str
    context_documents: List[str]
    query_intent: str
    group_context: Dict[str, Any]
    timestamp: datetime
    total_exchanges: int
    group_session_id: str

@dataclass
class QueryComplexityAnalysis:
    """Enhanced query complexity analysis for multi-level routing"""
    level: int
    level_name: str
    model_name: str
    max_tokens: int
    estimated_cost: float
    reasoning: str
    confidence: float
    expected_length: str
    complexity_level: str
    query_type: str
    academic_type: str
    key_concepts: List[str]
    sub_questions: List[str]
    requires_deep_search: bool
    requires_multi_round: bool
    requires_long_response: bool
    context_aware: bool
    force_detailed_response: bool = False
    explicit_detail_request: bool = False

@dataclass
class DocumentCategory:
    """Enhanced document categorization"""
    primary_subject: str
    secondary_subjects: List[str]
    document_type: str
    academic_level: str
    content_categories: List[str]
    key_topics: List[str]
    confidence_score: float
    language: str
    complexity_level: str

@dataclass
class EnhancedSearchResult:
    """Enhanced search result with comprehensive scoring"""
    document: Document
    relevance_score: float
    semantic_score: float
    keyword_score: float
    context_score: float
    category_match_score: float
    final_score: float
    search_method: str
    chunk_type: str
    source_document: str
    content_preview: str

@dataclass
class ProcessingResult:
    """Document processing result with enhanced metadata"""
    success: bool
    filename: str
    chunks_added: int
    document_type: str
    processing_time: float
    extraction_methods: List[str]
    metadata: Dict[str, Any]
    error: Optional[str] = None
    parallel_processing_used: bool = False
    thread_count: int = 1
    academic_content_detected: bool = False

@dataclass
class LargeDocumentProgress:
    """Progress tracking for large document processing"""
    total_chunks: int
    processed_chunks: int
    current_batch: int
    total_batches: int
    embedding_progress: float
    stage: str
    errors: List[str]

class UltraRobustAcademicDocumentProcessor:
    """
    Ultra-Robust Academic Document Processor with GPT-4.1 Series:
    - OPTIMIZED: Concise simple responses & comprehensive detailed explanations
    - NEW: GPT-4.1 series models with 1M token context support
    - Advanced categorization and subject detection
    - Multi-model query routing with conversation memory
    - Dynamic document size handling with massive context windows
    - Large document processing with enhanced batching
    - Adaptive response length management
    """

    def __init__(self, max_workers: int = None, enable_parallel: bool = True):
        """Initialize the ultra-robust academic processor with GPT-4.1 series"""
        

        self.vectorstore = None
        # Threading and performance configuration
        self.max_workers = max_workers or min(32, (os.cpu_count() or 1) + 4)
        self.enable_parallel = enable_parallel
        self.processing_lock = threading.Lock()
        self.conversation_lock = threading.Lock()

        # Initialize core AI components
        self.openai_api_key = os.getenv('OPENAI_API_KEY')
        if not self.openai_api_key:
            raise ValueError("OPENAI_API_KEY not found in environment variables")

        # Enhanced embeddings with higher dimensions for better semantic understanding
        self.embeddings = OpenAIEmbeddings(
            model="text-embedding-3-large",
            openai_api_key=self.openai_api_key,
            dimensions=3072
        )

        # Large document processing configuration - Enhanced for GPT-4.1 series
        self.embedding_batch_size = 300
        self.max_tokens_per_batch = 280000
        self.large_doc_threshold = 500000
        
        # Token encoder for cost calculation and batching
        self.token_encoder = tiktoken.encoding_for_model("gpt-4o-mini")

        # NEW: GPT-4.1 Series Model Configuration with 1M Token Context
        self.model_configs = {
            1: {  # Simple queries - ULTRA CONCISE with GPT-4.1 Nano
                "model_name": "gpt-4.1-nano",
                "input_cost_per_1m": 0.10,
                "output_cost_per_1m": 0.40,
                "max_tokens": 100,
                "max_context_tokens": 1000000,
                "temperature": 0.0,
                "description": "Ultra-fast GPT-4.1 Nano for simple factual queries",
                "use_cases": ["definitions", "dates", "simple facts", "yes/no questions"]
            },
            2: {  # Basic explanations - MODERATE with GPT-4.1 Mini
                "model_name": "gpt-4.1-mini",
                "input_cost_per_1m": 0.15,
                "output_cost_per_1m": 0.60,
                "max_tokens": 300,
                "max_context_tokens": 1000000,
                "temperature": 0.1,
                "description": "Balanced GPT-4.1 Mini for basic explanations",
                "use_cases": ["basic explanations", "simple how-to", "comparisons"]
            },
            3: {  # Detailed explanations - COMPREHENSIVE with GPT-4.1
                "model_name": "gpt-4.1",
                "input_cost_per_1m": 2.50,
                "output_cost_per_1m": 10.00,
                "max_tokens": 2000,
                "max_context_tokens": 1000000,
                "temperature": 0.2,
                "description": "Comprehensive GPT-4.1 for detailed explanations",
                "use_cases": ["detailed explanations", "examples", "procedures", "concepts"]
            },
            4: {  # Complex analysis - VERY COMPREHENSIVE with GPT-4.1
                "model_name": "gpt-4.1",
                "input_cost_per_1m": 2.50,
                "output_cost_per_1m": 10.00,
                "max_tokens": 3000,
                "max_context_tokens": 1000000,
                "temperature": 0.3,
                "description": "Advanced GPT-4.1 reasoning for complex academic queries",
                "use_cases": ["complex analysis", "multi-step reasoning", "academic research"]
            },
            5: {  # Expert synthesis - MAXIMUM COMPREHENSIVE with GPT-4.1
                "model_name": "gpt-4.1",
                "input_cost_per_1m": 2.50,
                "output_cost_per_1m": 10.00,
                "max_tokens": 5000,
                "max_context_tokens": 1000000,
                "temperature": 0.2,
                "description": "Expert-level GPT-4.1 comprehensive synthesis",
                "use_cases": ["synthesis", "expert analysis", "comprehensive research", "thesis-level work"]
            }
        }

        # ENHANCED: More context for GPT-4.1 series
        self.context_length_configs = {
            1: 3000,
            2: 5000,
            3: 10000,
            4: 20000,
            5: 50000
        }

        # OPTIMIZED: Enhanced complexity patterns with better detail detection
        self.complexity_patterns = {
            1: {  # Simple - Direct facts VERY STRICT FOR CONCISENESS
                "indicators": ["what is", "when is", "where is", "who is", "define", "name", "date", "time"],
                "negative_indicators": ["explain", "please explain", "with example", "detailed", "comprehensive", "how does", "why does", "elaborate"],
                "max_words": 5,
                "patterns": [r'\b(?:what|when|where|who)\s+is\b', r'\bdefine\s+\w+$', r'\bname\s+(?:the|a|an)\b'],
                "response_style": "ultra_concise",
                "complexity_boost": 0
            },
            2: {  # Basic - CAREFUL NOT TO TRIGGER FOR DETAILED REQUESTS
                "indicators": ["how", "why", "list", "summarize", "overview", "show"],
                "negative_indicators": ["please explain", "explain with example", "detailed explanation", "comprehensive", "elaborate", "complete concept"],
                "max_words": 10,
                "patterns": [r'\bhow\s+(?:to|does)\b', r'\bwhy\s+(?:is|does)\b', r'\blist\s+\w+$'],
                "response_style": "basic_explanation",
                "complexity_boost": 1
            },
            3: {  # Detailed explanations - PRIMARY TARGET FOR EXPLANATION REQUESTS
                "indicators": ["explain", "please explain", "explain this", "explain with example", "detailed explanation", "with example", "concept"],
                "negative_indicators": [],
                "max_words": 25,
                "patterns": [r'\bplease\s+explain\b', r'\bexplain\s+(?:this|with|the)\b', r'\bwith\s+example\b', r'\bdetailed\s+explanation\b', r'\bcomplete\s+concept\b'],
                "response_style": "detailed_comprehensive",
                "complexity_boost": 3
            },
            4: {  # Complex analysis - FOR DEEP ANALYSIS REQUESTS
                "indicators": ["analyze", "evaluate", "critical analysis", "comprehensive analysis", "in-depth", "thorough analysis"],
                "negative_indicators": [],
                "max_words": 35,
                "patterns": [r'\banalyze\s+(?:this|the)\b', r'\bevaluate\s+(?:this|the)\b', r'\bcomprehensive\s+analysis\b', r'\bin-depth\b'],
                "response_style": "complex_analytical",
                "complexity_boost": 4
            },
            5: {  # Expert synthesis - FOR MAXIMUM DETAIL REQUESTS
                "indicators": ["synthesize", "comprehensive analysis", "expert analysis", "complete explanation", "everything about"],
                "negative_indicators": [],
                "max_words": 100,
                "patterns": [r'\bsynthesize\b', r'\bcomprehensive\s+analysis\b', r'\bexpert\s+analysis\b', r'\beverything\s+about\b'],
                "response_style": "expert_comprehensive",
                "complexity_boost": 5
            }
        }

        # NEW: Explicit detail request patterns
        self.explicit_detail_patterns = [
            r'\bplease\s+explain\b',
            r'\bexplain\s+(?:this|with|the|complete|detailed?)\b',
            r'\bwith\s+example\b',
            r'\bdetailed?\s+explanation\b',
            r'\bcomplete\s+concept\b',
            r'\belaborate\s+(?:on|this)\b',
            r'\bcomprehensive\s+(?:explanation|analysis)\b',
            r'\bin-depth\s+(?:explanation|analysis)\b',
            r'\bthorough\s+(?:explanation|analysis)\b',
            r'\bexplain\s+fully\b',
            r'\btell\s+me\s+everything\b',
            r'\bwhile\s+solving\b',
            r'\beach\s+step\b'
        ]

        # Conversation Memory System
        self.conversation_contexts = {}
        self.conversation_memory_limit = 20

        self.group_conversation_contexts = {}
        self.group_conversation_memory_limit = 30  # More memory for group discussions
        self.group_conversation_lock = threading.Lock()

        # Context Relevance Thresholds
        self.relevance_thresholds = {
            "high_relevance": 0.75,
            "medium_relevance": 0.45,
            "low_relevance": 0.25,
            "minimal_relevance": 0.15
        }

        # Document processing and storage
        self.vector_store_path = "vectorstore"
        os.makedirs(self.vector_store_path, exist_ok=True)

        # Enhanced metadata with categorization
        self.document_metadata = {
            'processed_files': [],
            'last_update': '',
            'total_chunks': 0,
            'file_info': {},
            'categories': {},
            'subjects': {},
            'academic_levels': {}
        }

        # Advanced search components
        self.tfidf_vectorizer = TfidfVectorizer(
            max_features=15000,
            stop_words='english',
            ngram_range=(1, 3),
            min_df=1,
            max_df=0.85
        )
        self.tfidf_matrix = None
        self.tfidf_chunks = []
        self.tfidf_lock = threading.Lock()

        # Subject domains
        self.subject_domains = {
            "computer_science": {
                "keywords": ["algorithm", "programming", "software", "computer", "data structure", "database", "network", "AI", "machine learning", "cybersecurity", "web development", "mobile app"],
                "patterns": [r'\b(?:java|python|c\+\+|javascript|html|css|sql|nosql)\b', r'\bdata\s+structure\b', r'\bmachine\s+learning\b'],
                "subcategories": ["programming", "algorithms", "databases", "networks", "AI", "cybersecurity", "web_tech", "mobile_dev"]
            },
            "mathematics": {
                "keywords": ["calculus", "algebra", "geometry", "statistics", "probability", "linear algebra", "differential", "integral", "matrix", "equation", "theorem", "proof", "neural", "network", "learning"],
                "patterns": [r'\b(?:sin|cos|tan|log|ln|∫|∑|∆|π|θ|α|β|γ)\b', r'\bequation\b', r'\btheorem\b', r'\bproof\b', r'\bneural\s+network\b'],
                "subcategories": ["calculus", "algebra", "geometry", "statistics", "discrete_math", "linear_algebra", "neural_networks"]
            },
            "physics": {
                "keywords": ["force", "energy", "motion", "wave", "quantum", "relativity", "mechanics", "thermodynamics", "electromagnetism", "optics", "particle"],
                "patterns": [r'\b(?:newton|einstein|planck|maxwell)\b', r'\b(?:joule|watt|volt|amp|ohm)\b'],
                "subcategories": ["mechanics", "thermodynamics", "electromagnetism", "quantum_physics", "optics", "relativity"]
            },
            "chemistry": {
                "keywords": ["molecule", "atom", "reaction", "bond", "element", "compound", "organic", "inorganic", "periodic", "chemical", "catalyst", "acid", "base"],
                "patterns": [r'\b(?:H2O|CO2|NaCl|CH4|C6H12O6)\b', r'\b(?:acid|base|salt)\b', r'\bmolecular\s+formula\b'],
                "subcategories": ["organic_chemistry", "inorganic_chemistry", "physical_chemistry", "biochemistry"]
            },
            "biology": {
                "keywords": ["cell", "DNA", "RNA", "protein", "enzyme", "gene", "evolution", "ecosystem", "organism", "photosynthesis", "respiration", "genetics"],
                "patterns": [r'\b(?:DNA|RNA|ATP|NADH)\b', r'\bcell\s+division\b', r'\bnatural\s+selection\b'],
                "subcategories": ["cell_biology", "genetics", "ecology", "evolution", "biochemistry", "anatomy"]
            },
            "business": {
                "keywords": ["management", "marketing", "finance", "accounting", "economics", "strategy", "leadership", "operations", "supply chain", "human resources"],
                "patterns": [r'\b(?:ROI|NPV|IRR|SWOT|CEO|CFO)\b', r'\bbalance\s+sheet\b', r'\bmarket\s+analysis\b'],
                "subcategories": ["management", "marketing", "finance", "accounting", "economics", "entrepreneurship"]
            },
            "engineering": {
                "keywords": ["design", "manufacturing", "materials", "mechanical", "electrical", "civil", "structural", "circuit", "control systems", "engineering"],
                "patterns": [r'\b(?:CAD|CAM|FEA|CFD)\b', r'\bcontrol\s+system\b', r'\bstructural\s+analysis\b'],
                "subcategories": ["mechanical", "electrical", "civil", "chemical", "industrial", "software"]
            },
            "literature": {
                "keywords": ["novel", "poetry", "drama", "literature", "author", "character", "plot", "theme", "metaphor", "symbolism", "narrative"],
                "patterns": [r'\b(?:shakespeare|dickens|austen|hemingway)\b', r'\bliterary\s+analysis\b'],
                "subcategories": ["classic_literature", "modern_literature", "poetry", "drama", "literary_criticism"]
            },
            "history": {
                "keywords": ["war", "revolution", "empire", "civilization", "ancient", "medieval", "modern", "historical", "timeline", "dynasty"],
                "patterns": [r'\b(?:BC|AD|BCE|CE)\b', r'\b(?:war|revolution|empire)\b'],
                "subcategories": ["ancient_history", "medieval_history", "modern_history", "world_history", "political_history"]
            },
            "general_studies": {
                "keywords": ["education", "academic", "study", "course", "curriculum", "syllabus", "exam", "test", "assignment"],
                "patterns": [r'\b(?:exam|test|assignment|quiz)\b', r'\bsyllabus\b'],
                "subcategories": ["general_education", "study_materials", "exam_preparation"]
            }
        }

        # Document processors
        self.document_processors = {
            '.pdf': self._process_pdf_advanced,
            '.docx': self._process_docx_advanced,
            '.doc': self._process_docx_advanced,
            '.xlsx': self._process_excel_advanced,
            '.xls': self._process_excel_advanced,
            '.pptx': self._process_powerpoint_advanced,
            '.ppt': self._process_powerpoint_advanced,
            '.txt': self._process_text_advanced,
            '.html': self._process_html_advanced,
            '.htm': self._process_html_advanced
        }

        # Enhanced routing statistics
        self.routing_stats = {
            "total_queries": 0,
            "level_distribution": {1: 0, 2: 0, 3: 0, 4: 0, 5: 0},
            "total_cost": 0.0,
            "total_tokens": {"input": 0, "output": 0},
            "avg_response_times": {1: 0.0, 2: 0.0, 3: 0.0, 4: 0.0, 5: 0.0},
            "fallback_count": 0,
            "context_hits": 0,
            "follow_up_queries": 0,
            "explicit_detail_requests": 0,
            "gpt41_nano_queries": 0,
            "gpt41_mini_queries": 0,
            "gpt41_full_queries": 0
        }

        # Thread pools for enhanced performance
        self.executor = ThreadPoolExecutor(max_workers=self.max_workers)
        self.io_executor = ThreadPoolExecutor(max_workers=min(8, self.max_workers))

        # Initialize the system
        self._initialize_vectorstore()

    # ENHANCED: Query complexity analysis optimized for GPT-4.1 capabilities

    def _analyze_advanced_query_complexity(self, query: str, context: Optional[ConversationContext] = None) -> QueryComplexityAnalysis:
        """
        ENHANCED: Query complexity analysis optimized for GPT-4.1 series capabilities
        """
        query_lower = query.lower().strip()
        original_query = query.strip()
        
        # ENHANCED: Explicit detail request detection (HIGHEST PRIORITY)
        explicit_detail_request = any(re.search(pattern, query_lower) for pattern in self.explicit_detail_patterns)
        
        # Enhanced follow-up detection with immediate context awareness
        is_follow_up = False
        context_aware = False
        immediate_previous_context = ""
        
        if context and context.conversation_history:
            follow_up_indicators = ["continue", "more", "elaborate", "full", "complete", "detailed", "expand", "tell me more", "further", "additionally"]
            is_follow_up = any(indicator in query_lower for indicator in follow_up_indicators)
            context_aware = True
            
            if context.conversation_history:
                last_exchange = context.conversation_history[-1]
                immediate_previous_context = last_exchange.get('query', '')
                
                if any(term in immediate_previous_context.lower() for term in ['analyze', 'compare', 'evaluate', 'explain', 'detailed']):
                    is_follow_up = True

        # ENHANCED: Multi-factor complexity scoring optimized for GPT-4.1
        level_scores = {}
        
        for level in range(1, 6):
            score = 0.0
            patterns = self.complexity_patterns[level]
            
            # ENHANCED: Weighted positive indicators
            for indicator in patterns["indicators"]:
                if indicator in query_lower:
                    if query_lower.startswith(indicator) or f" {indicator}" in query_lower:
                        score += 3.5
                    else:
                        score += 2.5
            
            # ENHANCED: Regex pattern matching with higher weights
            for pattern in patterns["patterns"]:
                matches = len(re.findall(pattern, query_lower))
                score += matches * 3.0
            
            # ENHANCED: Negative indicators (stronger penalty for simple queries)
            for neg_indicator in patterns["negative_indicators"]:
                if neg_indicator in query_lower:
                    if level <= 2:
                        score -= 5.0
                    else:
                        score -= 2.0
            
            # ENHANCED: Word count analysis (more strict for simple queries)
            word_count = len(original_query.split())
            if word_count <= patterns["max_words"]:
                score += 1.5
            elif word_count > patterns["max_words"] * 2:
                if level <= 2:
                    score -= 3.0
            
            # NEW: Complexity boost based on academic indicators
            academic_complexity_markers = {
                1: ["quick", "short", "brief", "just"],
                2: ["show", "tell", "how"],
                3: ["explain", "describe", "example", "concept", "detailed"],
                4: ["analyze", "evaluate", "critical", "comprehensive", "thorough"],
                5: ["synthesize", "expert", "research", "complete analysis"]
            }
            
            for marker in academic_complexity_markers.get(level, []):
                if marker in query_lower:
                    score += 2.0 + patterns.get("complexity_boost", 0)
            
            # NEW: Enhanced question complexity indicators
            complexity_indicators = [
                (r'\bwhy\s+(?:does|is|are)', 2),
                (r'\bhow\s+(?:does|can|do)', 1),
                (r'\bwhat\s+are\s+the\s+(?:steps|differences|similarities)', 2),
                (r'\b(?:compare|contrast)\s+and\b', 3),
                (r'\b(?:analyze|evaluate|assess|critique)\b', 3),
                (r'\b(?:synthesize|comprehensive)\b', 4)
            ]
            
            for pattern, boost in complexity_indicators:
                if re.search(pattern, query_lower):
                    score += boost
            
            # CRITICAL: Explicit detail request override for GPT-4.1
            if explicit_detail_request and level >= 3:
                score += 15.0
            elif explicit_detail_request and level <= 2:
                score -= 7.0
            
            level_scores[level] = max(score, 0.0)
        
        # ENHANCED: Better level determination optimized for GPT-4.1
        if explicit_detail_request:
            # FORCE appropriate level for detail requests with GPT-4.1
            if any(pattern in query_lower for pattern in ['comprehensive', 'complete concept', 'everything', 'in-depth']):
                best_level = 5
            elif any(pattern in query_lower for pattern in ['detailed', 'thorough', 'elaborate']):
                best_level = 4
            else:
                best_level = 3
            
        elif all(score == 0 for score in level_scores.values()):
            # STRICT: Default assignment favoring conciseness with GPT-4.1 Nano
            if any(word in query_lower for word in ['what is', 'when is', 'define']):
                best_level = 1
            elif any(word in query_lower for word in ['how', 'why']) and not any(word in query_lower for word in ['explain', 'detailed']):
                best_level = 2
            else:
                best_level = 2
        else:
            best_level = max(level_scores, key=level_scores.get)
            
            # ENHANCED: Level validation with explicit detail consideration
            if level_scores[best_level] < 2.0 and best_level > 2 and not explicit_detail_request:
                sorted_levels = sorted(level_scores.items(), key=lambda x: x[1], reverse=True)
                if len(sorted_levels) > 1:
                    best_level = min(sorted_levels[1][0], 2)
        
        # Enhanced follow-up query handling
        if is_follow_up and context:
            previous_complexity = context.follow_up_context.get("complexity_level", 2)
            
            if any(indicator in query_lower for indicator in ['more detail', 'elaborate', 'expand', 'comprehensive']):
                best_level = min(previous_complexity + 1, 5)
            elif 'continue' in query_lower:
                best_level = max(previous_complexity, best_level)
            
            if immediate_previous_context:
                if any(term in immediate_previous_context.lower() for term in ['analyze', 'compare', 'evaluate']):
                    best_level = max(best_level, 3)
                elif any(term in immediate_previous_context.lower() for term in ['synthesize', 'comprehensive']):
                    best_level = max(best_level, 4)

        # ENHANCED: Response length expectation for GPT-4.1
        requires_long_response = explicit_detail_request or best_level >= 3 or any(indicator in query_lower for indicator in [
            "detailed", "comprehensive", "complete", "elaborate", "in-depth", "thorough",
            "explain fully", "tell me everything", "all about"
        ])

        # ENHANCED: Force detailed response detection for GPT-4.1
        force_detailed_response = explicit_detail_request or any(pattern in query_lower for pattern in [
            "detailed explanation", "comprehensive analysis", "tell me everything",
            "complete concept", "thorough explanation", "explain fully",
            "while solving", "each step", "step by step"
        ])

        # Get model configuration
        model_config = self.model_configs[best_level]
        
        # ENHANCED: Cost estimation for GPT-4.1 series
        estimated_input_tokens = len(self.token_encoder.encode(query)) + 3000
        estimated_output_tokens = model_config["max_tokens"]
        
        estimated_cost = (
            (estimated_input_tokens / 1_000_000) * model_config["input_cost_per_1m"] +
            (estimated_output_tokens / 1_000_000) * model_config["output_cost_per_1m"]
        )
        
        # ENHANCED: Reasoning generation with GPT-4.1 context
        reasoning_parts = [f"Selected Level {best_level} - {model_config['model_name']} ({model_config['description']})"]
        
        if explicit_detail_request:
            reasoning_parts.append(f"EXPLICIT DETAIL REQUEST DETECTED - Using {model_config['model_name']} with comprehensive response mode")
        
        if level_scores[best_level] > 0:
            top_indicators = []
            for indicator in self.complexity_patterns[best_level]["indicators"]:
                if indicator in query_lower:
                    top_indicators.append(indicator)
            if top_indicators:
                reasoning_parts.append(f"• Detected patterns: {', '.join(top_indicators[:3])}")
        
        if is_follow_up:
            reasoning_parts.append("• Follow-up query detected - maintaining context continuity")
        
        if requires_long_response:
            reasoning_parts.append(f"• Comprehensive response required - leveraging {model_config['model_name']} capabilities")
            
        if force_detailed_response:
            reasoning_parts.append(f"• Forcing detailed response with {model_config['model_name']}")
            
        if immediate_previous_context:
            reasoning_parts.append(f"• Previous context: '{immediate_previous_context[:50]}...'")
        
        reasoning = "\n".join(reasoning_parts)
        
        # ENHANCED: Level names with GPT-4.1 models
        level_names = {
            1: "Simple Query - GPT-4.1 Nano (Concise)",
            2: "Basic Analysis - GPT-4.1 Mini (Moderate)",
            3: "Detailed Explanation - GPT-4.1 (Comprehensive)",
            4: "Complex Analysis - GPT-4.1 (Very Comprehensive)", 
            5: "Expert Synthesis - GPT-4.1 (Maximum Detail)"
        }
        
        # Academic type and query analysis
        academic_type = self._detect_advanced_academic_type(query_lower)
        query_type = self._detect_enhanced_query_type(query_lower)
        
        return QueryComplexityAnalysis(
            level=best_level,
            level_name=level_names[best_level],
            model_name=model_config["model_name"],
            max_tokens=model_config["max_tokens"],
            estimated_cost=estimated_cost,
            reasoning=reasoning,
            confidence=min(level_scores[best_level] / 5.0, 0.95),
            expected_length="comprehensive" if requires_long_response else "concise",
            complexity_level=self.complexity_patterns[best_level]["response_style"],
            query_type=query_type,
            academic_type=academic_type,
            key_concepts=self._extract_enhanced_key_concepts(query),
            sub_questions=[],
            requires_deep_search=best_level >= 3,
            requires_multi_round=best_level >= 4,
            requires_long_response=requires_long_response,
            context_aware=context_aware,
            force_detailed_response=force_detailed_response,
            explicit_detail_request=explicit_detail_request
        )

    def _create_enhanced_prompt(self, query: str, context: str, query_analysis: QueryComplexityAnalysis, 
                              conversation_context: ConversationContext) -> str:
        """ENHANCED: Create prompts optimized for GPT-4.1 series with massive context"""
        
        # ENHANCED: Base system prompt for GPT-4.1 series
        system_prompt = f"""You are an advanced academic AI assistant powered by {query_analysis.model_name} with access to massive context windows. Your role is to provide accurate, helpful, and contextually appropriate responses based on the provided academic materials.

GPT-4.1 ANALYSIS CONTEXT:
- Model: {query_analysis.model_name}
- Query Complexity Level: {query_analysis.level} ({query_analysis.level_name})
- Expected Response Length: {query_analysis.expected_length}
- Explicit Detail Request: {query_analysis.explicit_detail_request}
- Academic Type: {query_analysis.academic_type}
- Query Type: {query_analysis.query_type}
- Key Concepts: {', '.join(query_analysis.key_concepts[:5])}
- Context Capacity: {self.model_configs[query_analysis.level]['max_context_tokens']:,} tokens
"""

        # ENHANCED: Add conversation context optimized for GPT-4.1
        if conversation_context.conversation_history:
            recent_history = conversation_context.conversation_history[-5:]
            system_prompt += f"""
ENHANCED CONVERSATION CONTEXT (GPT-4.1 Memory):
Total exchanges: {len(conversation_context.conversation_history)}
"""
            if recent_history:
                last_exchange = recent_history[-1]
                system_prompt += f"""
IMMEDIATE PREVIOUS CONTEXT (High Priority):
Previous Question: {last_exchange['query']}
Previous Response: {last_exchange['response'][:300]}...  
"""
                # ENHANCED: More conversation history for GPT-4.1
                if len(recent_history) > 1:
                    system_prompt += "Recent Conversation History:\n"
                    for i, exchange in enumerate(recent_history[:-1], 1):
                        system_prompt += f"Exchange {i}: Q: {exchange['query'][:100]}... A: {exchange['response'][:100]}...\n"

        # ENHANCED: Level-specific instructions optimized for GPT-4.1 capabilities
        level_instructions = {
            1: {  # GPT-4.1 Nano - ULTRA-CONCISE
                "style": "Provide an ultra-direct, concise answer using GPT-4.1 Nano's efficiency. Give ONLY what the user asked for in 1 sentence maximum.",
                "format": "Use simple, clear language. Bold key facts with **text**.",
                "length": "Target: 10-40 words. Be extremely concise and precise.",
                "enforcement": "CRITICAL: GPT-4.1 Nano mode - Keep response VERY brief. User wants just the answer, nothing more."
            },
            2: {  # GPT-4.1 Mini - MODERATE
                "style": "Provide a clear, balanced explanation using GPT-4.1 Mini's capabilities in 2-3 sentences maximum.",
                "format": "Use educational tone. Bold important terms with **text**.",
                "length": "Target: 40-120 words. Keep it focused and to the point.",
                "enforcement": "IMPORTANT: GPT-4.1 Mini mode - Be concise but complete. Don't over-explain unless asked."
            },
            3: {  # GPT-4.1 Full - COMPREHENSIVE
                "style": "Leverage GPT-4.1's comprehensive capabilities to provide detailed explanation with examples and context.",
                "format": "Use bullet points, examples, and clear structure. Bold key concepts. Utilize the 1M token context effectively.",
                "length": "Target: 300-800 words. Be thorough and comprehensive with rich detail.",
                "enforcement": "CRITICAL: GPT-4.1 Full mode - User requested detailed explanation. Provide comprehensive response with examples and thorough coverage using the massive context available."
            },
            4: {  # GPT-4.1 Full - VERY COMPREHENSIVE
                "style": "Use GPT-4.1's advanced reasoning for in-depth analysis with step-by-step reasoning and detailed examples.",
                "format": "Use numbered lists, headings, examples. Show clear logical progression. Leverage the 1M token context for comprehensive analysis.",
                "length": "Target: 600-1500 words. Be very thorough and analytical.",
                "enforcement": "CRITICAL: GPT-4.1 Full mode - This is a complex query requiring substantial detailed response with comprehensive analysis using the full context capabilities."
            },
            5: {  # GPT-4.1 Full - MAXIMUM COMPREHENSIVE
                "style": "Utilize GPT-4.1's maximum capabilities for expert-level comprehensive analysis with complete coverage.",
                "format": "Structure: Overview → Detailed Analysis → Examples → Applications → Synthesis → Conclusion. Use headings and comprehensive formatting.",
                "length": "Target: 1000-3000+ words. Be extremely thorough and comprehensive.",
                "enforcement": "CRITICAL: GPT-4.1 Full Expert mode - Provide maximum comprehensive response. User expects expert-level detailed explanation using the full 1M token context."
            }
        }
        
        current_instructions = level_instructions[query_analysis.level]
        
        # CRITICAL: Override for explicit detail requests with GPT-4.1
        if query_analysis.explicit_detail_request and query_analysis.level < 3:
            current_instructions = level_instructions[3]
            current_instructions["enforcement"] = f"OVERRIDE: User explicitly requested detailed explanation. Using GPT-4.1's comprehensive capabilities to ignore typical length constraints and provide detailed response."

        # ENHANCED: Query type specific instructions for GPT-4.1
        query_type_instructions = {
            "definition": "Leverage GPT-4.1's knowledge to provide clear, precise definition. For simple requests, be concise. For detailed requests, include comprehensive examples and context.",
            "explanation": "Use GPT-4.1's reasoning capabilities. For simple 'how' questions, be direct. For 'please explain' requests, be comprehensive with examples and detailed analysis.",
            "comparison": "Structure comparison using GPT-4.1's analytical capabilities. Be concise for simple comparisons, comprehensive for detailed analysis with rich context.",
            "analysis": "Leverage GPT-4.1's advanced reasoning for structured analysis. Match depth to user's request level using the massive context available.",
            "follow_up": "Build on previous context using GPT-4.1's memory capabilities. Match the detail level to user's follow-up request.",
            "procedural": "Use GPT-4.1's step-by-step reasoning. For simple 'how to' be direct. For detailed explanation requests, provide comprehensive step-by-step guide with examples."
        }
        
        # ENHANCED: Construct optimized prompt for GPT-4.1
        final_prompt = f"""{system_prompt}

GPT-4.1 RESPONSE REQUIREMENTS:
- Model Capability: {query_analysis.model_name}
- Style: {current_instructions['style']}
- Format: {current_instructions['format']}
- Length: {current_instructions['length']}
- Enforcement: {current_instructions['enforcement']}
- Special Instructions: {query_type_instructions.get(query_analysis.query_type, 'Match response length to user request using GPT-4.1 capabilities.')}

CRITICAL GPT-4.1 RESPONSE RULES:
1. Answer ONLY based on the provided context - do not add external information
2. If information is not available in the context, clearly state this
3. Use **bold formatting** for key terms, numbers, and important concepts, dont give raw text expect markdowns
4. LEVERAGE GPT-4.1's MASSIVE CONTEXT WINDOW:
   - Utilize up to 1M tokens of context for comprehensive responses
   - Cross-reference information across large documents
   - Maintain coherence across extensive content
5. MATCH RESPONSE LENGTH TO USER'S REQUEST EXACTLY:
   - Simple questions = Concise answers with GPT-4.1 Nano efficiency
   - "Please explain" or "with example" = Comprehensive detailed responses with GPT-4.1 Full
   - "Detailed explanation" = Very comprehensive responses using full GPT-4.1 capabilities
6. If this is a follow-up query, acknowledge and build upon previous context using GPT-4.1's memory
7. For Level 1-2: Use GPT-4.1 Nano/Mini efficiency - Be concise, direct, focused
8. For Level 3+: Use GPT-4.1 Full capabilities - Be comprehensive, detailed, thorough with examples
9. NEVER provide long responses for simple questions
10. NEVER provide brief responses when user explicitly asks for detailed explanation
11. LEVERAGE THE 1M TOKEN CONTEXT for comprehensive analysis when appropriate
12. do not mention sources or context where data is retrived

ACADEMIC CONTEXT (Utilize GPT-4.1's massive context capacity):
{context}

STUDENT QUERY: {query}

PROVIDE YOUR GPT-4.1 OPTIMIZED ACADEMIC RESPONSE (Remember to match length precisely to user's request level using appropriate GPT-4.1 model capabilities):"""
        
        return final_prompt

    def _adaptive_response_enhancement(self, query: str, response: str, query_analysis: QueryComplexityAnalysis, 
                                     context: str, llm: ChatOpenAI) -> str:
        """
        ENHANCED: Adaptive response enhancement optimized for GPT-4.1 series capabilities
        """
        query_lower = query.lower()
        current_length = len(response)
        
        # ENHANCED: Detection of user intent with GPT-4.1 awareness
        explicit_expansion_requests = [
            "please explain", "explain with example", "detailed explanation", 
            "elaborate", "comprehensive", "tell me more", "in detail",
            "thorough", "complete concept", "explain fully", "everything about"
        ]
        
        explicit_contraction_requests = [
            "brief", "short", "quick", "simple", "just", "only", "concise"
        ]
        
        wants_expansion = any(req in query_lower for req in explicit_expansion_requests)
        wants_contraction = any(req in query_lower for req in explicit_contraction_requests)
        
        # ENHANCED: Context availability assessment for GPT-4.1
        context_richness = len(context) / 1000
        available_info_score = min(context_richness, 5.0)
        
        # CRITICAL: Handle explicit detail requests with GPT-4.1
        if query_analysis.explicit_detail_request and current_length < 500 and available_info_score > 0.5:
            
            expansion_prompt = f"""The user has explicitly requested a detailed explanation. Leverage GPT-4.1's comprehensive capabilities and the massive context available to provide a much more detailed response.

ORIGINAL QUERY (requesting detailed explanation): {query}

RICH CONTEXT AVAILABLE (GPT-4.1 can process up to 1M tokens):
{context[:3000]}...

CURRENT BRIEF RESPONSE:
{response}

GPT-4.1 COMPREHENSIVE DETAILED RESPONSE: Use GPT-4.1's advanced capabilities to expand this significantly with:
- Detailed explanations leveraging the full context
- Multiple examples where applicable  
- Step-by-step breakdowns if relevant
- Comprehensive coverage of the topic using available information
- Clear structure and formatting optimized for readability
- Cross-references across the large context when relevant

The user explicitly wants detailed information, so use GPT-4.1's full capabilities to provide a thorough, comprehensive response."""

            try:
                enhanced_result = llm.invoke(expansion_prompt)
                if len(enhanced_result.content) > current_length:
                    return enhanced_result.content
            except Exception as e:
                pass
        
        # Handle explicit brevity requests with GPT-4.1 efficiency
        elif wants_contraction and current_length > 100:
            
            condensation_prompt = f"""The user requested a brief, concise response. Use GPT-4.1's efficiency to provide a condensed version that covers just the essential information.

ORIGINAL RESPONSE:
{response}

GPT-4.1 BRIEF RESPONSE: Provide a concise, direct answer focusing only on what the user specifically asked for, leveraging GPT-4.1's ability to distill information efficiently."""

            try:
                condensed_result = llm.invoke(condensation_prompt)
                if len(condensed_result.content) < current_length and len(condensed_result.content) > 30:
                    return condensed_result.content
            except Exception as e:
                pass
        
        # Auto-expand for Level 3+ queries with GPT-4.1 capabilities
        elif query_analysis.level >= 3 and current_length < 400 and available_info_score > 1.0:
            
            auto_expansion_prompt = f"""This is a Level {query_analysis.level} query requiring a more detailed response. Use GPT-4.1's comprehensive capabilities and the rich context available to provide a more thorough answer.

RICH CONTEXT (GPT-4.1 can leverage extensive context):
{context[:2000]}...

CURRENT RESPONSE:
{response}

GPT-4.1 ENHANCED DETAILED RESPONSE: Provide a more comprehensive response appropriate for this complexity level, using GPT-4.1's ability to synthesize information from extensive context with better explanations, examples, and detailed analysis."""

            try:
                auto_enhanced_result = llm.invoke(auto_expansion_prompt)
                if len(auto_enhanced_result.content) > current_length:
                    return auto_enhanced_result.content
            except Exception as e:
                pass
        
        # Ensure appropriate conciseness for Level 1-2 with GPT-4.1 efficiency
        elif query_analysis.level <= 2 and current_length > 200:
            
            concise_prompt = f"""This is a Level {query_analysis.level} query requiring a concise response. Use GPT-4.1's efficiency to provide a brief, direct answer.

CURRENT RESPONSE:
{response}

GPT-4.1 CONCISE RESPONSE: Provide a brief, direct answer that gives just what the user asked for without unnecessary elaboration, leveraging GPT-4.1's ability to be precise and efficient."""

            try:
                concise_result = llm.invoke(concise_prompt)
                if len(concise_result.content) < current_length and len(concise_result.content) > 30:
                    return concise_result.content
            except Exception as e:
                pass
        
        return response

    # ALL LARGE DOCUMENT PROCESSING METHODS (ENHANCED FOR GPT-4.1)
    
    def _count_tokens_in_chunks(self, chunks: List[Document]) -> int:
        """Count total tokens in a list of chunks"""
        total_tokens = 0
        for chunk in chunks:
            try:
                tokens = len(self.token_encoder.encode(chunk.page_content))
                total_tokens += tokens
            except Exception as e:
                total_tokens += len(chunk.page_content) // 4
        return total_tokens

    def _create_embedding_batches(self, chunks: List[Document]) -> List[List[Document]]:
        """ENHANCED: Create embedding batches optimized for GPT-4.1 processing"""
        batches = []
        current_batch = []
        current_tokens = 0
        
        for i, chunk in enumerate(chunks):
            try:
                chunk_tokens = len(self.token_encoder.encode(chunk.page_content))
            except Exception:
                chunk_tokens = len(chunk.page_content) // 4
            
            # ENHANCED: Larger batches for GPT-4.1 processing
            if (current_tokens + chunk_tokens > self.max_tokens_per_batch or 
                len(current_batch) >= self.embedding_batch_size):
                
                if current_batch:
                    batches.append(current_batch)
                
                current_batch = [chunk]
                current_tokens = chunk_tokens
            else:
                current_batch.append(chunk)
                current_tokens += chunk_tokens
        
        if current_batch:
            batches.append(current_batch)
        
        return batches

    def _process_embedding_batch(self, batch: List[Document], batch_num: int, total_batches: int) -> bool:
        """ENHANCED: Process embedding batch with GPT-4.1 optimization"""
        max_retries = 3
        retry_delay = 2.0
        
        for attempt in range(max_retries):
            try:
                
                total_tokens = self._count_tokens_in_chunks(batch)
                if total_tokens > self.max_tokens_per_batch:
                    mid_point = len(batch) // 2
                    if mid_point > 0:
                        batch1 = batch[:mid_point]
                        batch2 = batch[mid_point:]
                        
                        success1 = self._process_embedding_batch(batch1, f"{batch_num}a", total_batches)
                        success2 = self._process_embedding_batch(batch2, f"{batch_num}b", total_batches)
                        
                        return success1 and success2
                    else:
                        return False
                
                with self.processing_lock:
                    if self.vectorstore is None:
                        self.vectorstore = FAISS.from_documents(batch, self.embeddings)
                    else:
                        self.vectorstore.add_documents(batch)
                
                return True
                
            except Exception as e:
                error_msg = str(e)
                if "max_tokens_per_request" in error_msg or "400" in error_msg:
                    if len(batch) > 1:
                        mid_point = len(batch) // 2
                        batch1 = batch[:mid_point]
                        batch2 = batch[mid_point:]
                        
                        success1 = self._process_embedding_batch(batch1, f"{batch_num}a", total_batches)
                        success2 = self._process_embedding_batch(batch2, f"{batch_num}b", total_batches)
                        
                        return success1 and success2
                    else:
                        return False
                else:
                    if attempt < max_retries - 1:
                        time.sleep(retry_delay)
                        retry_delay *= 1.5
                    else:
                        return False
        
        return False

    def _process_large_document_embeddings(self, chunks: List[Document], filename: str) -> bool:
        """ENHANCED: Process embeddings for large documents with GPT-4.1 optimization"""
        
        total_tokens = self._count_tokens_in_chunks(chunks)
        
        if total_tokens <= self.max_tokens_per_batch and len(chunks) <= self.embedding_batch_size:
            try:
                with self.processing_lock:
                    if self.vectorstore is None:
                        self.vectorstore = FAISS.from_documents(chunks, self.embeddings)
                    else:
                        self.vectorstore.add_documents(chunks)
                return True
            except Exception as e:
                return False
        
        batches = self._create_embedding_batches(chunks)
        
        if not batches:
            return False
        
        successful_batches = 0
        failed_batches = 0
        
        for i, batch in enumerate(batches, 1):
            try:
                success = self._process_embedding_batch(batch, i, len(batches))
                if success:
                    successful_batches += 1
                else:
                    failed_batches += 1
                
                if i < len(batches):
                    time.sleep(0.5)
                    
            except Exception as e:
                failed_batches += 1
        
        try:
            with self.processing_lock:
                if self.vectorstore:
                    self.vectorstore.save_local(self.vector_store_path)
        except Exception as e:
            pass
        
        return successful_batches > 0

    # ALL OTHER HELPER METHODS (PRESERVED WITH GPT-4.1 ENHANCEMENTS)
    
    def _detect_advanced_academic_type(self, query_lower: str) -> str:
        """Enhanced academic type detection"""
        if any(term in query_lower for term in ['question', 'q.', 'marks', 'exam', 'test', 'quiz', 'assignment']):
            return 'question_paper'
        elif any(term in query_lower for term in ['answer', 'solution', 'solve', 'response', 'explanation']):
            return 'answer_script'
        elif any(term in query_lower for term in ['textbook', 'book', 'chapter', 'section', 'theory', 'concept']):
            return 'textbook_content'
        elif any(term in query_lower for term in ['research', 'paper', 'journal', 'article', 'study']):
            return 'research_material'
        elif any(term in query_lower for term in ['lecture', 'notes', 'slides', 'presentation']):
            return 'lecture_material'
        else:
            return 'general_academic'

    def _detect_enhanced_query_type(self, query_lower: str) -> str:
        """Enhanced query type detection with more categories"""
        if any(word in query_lower for word in ['define', 'definition', 'meaning', 'what is']):
            return 'definition'
        elif any(word in query_lower for word in ['explain', 'how does', 'how to', 'describe']):
            return 'explanation'
        elif any(word in query_lower for word in ['compare', 'contrast', 'difference', 'similarity', 'versus']):
            return 'comparison'
        elif any(word in query_lower for word in ['analyze', 'analysis', 'evaluate', 'assess', 'examine']):
            return 'analysis'
        elif any(word in query_lower for word in ['list', 'enumerate', 'name', 'identify']):
            return 'listing'
        elif any(word in query_lower for word in ['summarize', 'summary', 'overview', 'outline']):
            return 'summary'
        elif any(word in query_lower for word in ['solve', 'calculate', 'compute', 'find']):
            return 'problem_solving'
        elif any(word in query_lower for word in ['examples', 'example', 'instance', 'case']):
            return 'examples'
        elif any(word in query_lower for word in ['steps', 'procedure', 'process', 'method', 'algorithm']):
            return 'procedural'
        elif any(word in query_lower for word in ['continue', 'more', 'elaborate', 'expand', 'detailed']):
            return 'follow_up'
        else:
            return 'general_inquiry'

    def _extract_enhanced_key_concepts(self, query: str) -> List[str]:
        """Enhanced key concept extraction with domain-specific terms"""
        tokens = safe_word_tokenize(query.lower())
        key_concepts = []
        
        academic_terms = set([
            'algorithm', 'data', 'structure', 'database', 'network', 'security', 'machine', 'learning',
            'artificial', 'intelligence', 'programming', 'software', 'hardware', 'system',
            'mathematics', 'calculus', 'algebra', 'geometry', 'statistics', 'probability',
            'physics', 'chemistry', 'biology', 'engineering', 'business', 'management',
            'literature', 'history', 'philosophy', 'psychology', 'sociology', 'economics',
            'neural', 'networks', 'neurons', 'synapses', 'perceptron', 'backpropagation'
        ])
        
        for token in tokens:
            if (len(token) > 3 and
                token not in STOP_WORDS and
                token.isalpha()):
                lemmatized = safe_lemmatize(token)
                
                if token in academic_terms or lemmatized in academic_terms:
                    key_concepts.insert(0, lemmatized)
                else:
                    key_concepts.append(lemmatized)
        
        return key_concepts[:8]

    def _get_or_create_conversation_context(self, session_id: str) -> ConversationContext:
        """Get or create conversation context for a session"""
        with self.conversation_lock:
            if session_id not in self.conversation_contexts:
                self.conversation_contexts[session_id] = ConversationContext(
                    session_id=session_id,
                    conversation_history=[],
                    last_query="",
                    last_response="",
                    context_documents=[],
                    query_intent="",
                    follow_up_context={},
                    timestamp=datetime.now(),
                    total_exchanges=0
                )
            return self.conversation_contexts[session_id]

    def _update_conversation_context(self, session_id: str, query: str, response: str, 
                                   query_analysis: QueryComplexityAnalysis, sources: List[str]):
        """ENHANCED: Update conversation context with GPT-4.1 awareness"""
        with self.conversation_lock:
            context = self.conversation_contexts.get(session_id)
            if context:
                context.conversation_history.append({
                    "query": query,
                    "response": response,
                    "timestamp": datetime.now().isoformat(),
                    "complexity_level": query_analysis.level,
                    "sources": sources,
                    "query_type": query_analysis.query_type,
                    "response_length": len(response),
                    "explicit_detail_request": query_analysis.explicit_detail_request,
                    "model_used": query_analysis.model_name
                })
                
                if len(context.conversation_history) > self.conversation_memory_limit:
                    context.conversation_history = context.conversation_history[-self.conversation_memory_limit:]
                
                context.last_query = query
                context.last_response = response
                context.context_documents = sources
                context.query_intent = query_analysis.query_type
                context.follow_up_context = {
                    "complexity_level": query_analysis.level,
                    "academic_type": query_analysis.academic_type,
                    "key_concepts": query_analysis.key_concepts,
                    "response_length": len(response),
                    "detailed_response_given": len(response) > 300,
                    "explicit_detail_request": query_analysis.explicit_detail_request,
                    "gpt41_model_used": query_analysis.model_name
                }
                context.timestamp = datetime.now()
                context.total_exchanges += 1

    def _intelligent_context_relevance_scoring(self, query: str, documents: List[Document], 
                                             query_analysis: QueryComplexityAnalysis) -> List[EnhancedSearchResult]:
        """ENHANCED: Context relevance scoring optimized for GPT-4.1's massive context"""
        enhanced_results = []
        query_lower = query.lower()
        query_concepts = set(query_analysis.key_concepts)
        
        for doc in documents:
            content = doc.page_content.lower()
            metadata = doc.metadata
            
            try:
                content_words = set(safe_word_tokenize(content))
                query_words = set(safe_word_tokenize(query_lower))
                
                common_words = content_words.intersection(query_words)
                semantic_score = len(common_words) / (len(query_words) + 1)
            except Exception:
                semantic_score = 0.0
            
            keyword_score = 0.0
            query_words = query_lower.split()
            
            for word in query_words:
                if word in content:
                    count = content.count(word)
                    first_pos = content.find(word)
                    if first_pos != -1:
                        position_weight = 1.0 - (first_pos / len(content))
                        keyword_score += count * position_weight
            
            keyword_score = min(keyword_score / len(query_words), 1.0)
            
            context_score = 0.0
            
            if query_analysis.academic_type != 'general_academic':
                if query_analysis.academic_type == 'question_paper':
                    if any(pattern in content for pattern in ['question', 'q.', 'marks', 'exam']):
                        context_score += 0.3
                elif query_analysis.academic_type == 'answer_script':
                    if any(pattern in content for pattern in ['answer', 'solution', 'explanation']):
                        context_score += 0.3
                elif query_analysis.academic_type == 'textbook_content':
                    if any(pattern in content for pattern in ['chapter', 'section', 'theory']):
                        context_score += 0.3
            
            if query_analysis.query_type == 'definition':
                if any(pattern in content for pattern in ['definition', 'defined as', 'refers to']):
                    context_score += 0.2
            elif query_analysis.query_type == 'explanation':
                if any(pattern in content for pattern in ['explanation', 'because', 'due to', 'reason']):
                    context_score += 0.2
            elif query_analysis.query_type == 'comparison':
                if any(pattern in content for pattern in ['compare', 'versus', 'difference', 'similar']):
                    context_score += 0.2
            
            category_match_score = 0.0
            if 'category' in metadata:
                doc_category = metadata.get('category', {})
                if isinstance(doc_category, dict):
                    primary_subject = doc_category.get('primary_subject', '')
                    if primary_subject:
                        subject_keywords = self.subject_domains.get(primary_subject, {}).get('keywords', [])
                        matching_keywords = sum(1 for keyword in subject_keywords if keyword in query_lower)
                        category_match_score = min(matching_keywords / 5.0, 1.0)
            
            quality_indicators = ['example', 'definition', 'explanation', 'step', 'method', 'procedure']
            quality_score = sum(1 for indicator in quality_indicators if indicator in content) / len(quality_indicators)
            
            # ENHANCED: Weights optimized for GPT-4.1's context capabilities
            weights = {
                'semantic': 0.25,
                'keyword': 0.30,
                'context': 0.20,
                'category': 0.15,
                'quality': 0.10
            }
            
            final_score = (
                semantic_score * weights['semantic'] +
                keyword_score * weights['keyword'] +
                context_score * weights['context'] +
                category_match_score * weights['category'] +
                quality_score * weights['quality']
            )
            
            content_preview = content[:250] + "..." if len(content) > 250 else content
            
            enhanced_result = EnhancedSearchResult(
                document=doc,
                relevance_score=final_score,
                semantic_score=semantic_score,
                keyword_score=keyword_score,
                context_score=context_score,
                category_match_score=category_match_score,
                final_score=final_score,
                search_method="gpt41_intelligent_relevance",
                chunk_type=metadata.get('chunk_type', 'unknown'),
                source_document=metadata.get('source', 'unknown'),
                content_preview=content_preview
            )
            
            enhanced_results.append(enhanced_result)
        
        enhanced_results.sort(key=lambda x: x.final_score, reverse=True)
        return enhanced_results

    def generate_response(self, query: str, session_id: str = "default", user_context: Dict = None) -> Dict[str, Any]:
        """
        GPT-4.1 OPTIMIZED: Generate response with precise length control using GPT-4.1 series
        """
        start_time = time.time()
        
        try:
            
            # Get conversation context
            conversation_context = self._get_or_create_conversation_context(session_id)
            
            # GPT-4.1 OPTIMIZED: Enhanced complexity analysis
            query_analysis = self._analyze_advanced_query_complexity(query, conversation_context)
            
            # Check document availability
            if not self.vectorstore or self.vectorstore.index.ntotal <= 1:
                self._update_conversation_context(session_id, query, 
                                                "No documents available. Please upload documents first.",
                                                query_analysis, [])
                return {
                    "response": "I don't have any documents to reference. Please upload some academic materials first, and I'll be happy to help answer your questions based on that content! 📚",
                    "complexity_analysis": {
                        "level": query_analysis.level,
                        "level_name": query_analysis.level_name,
                        "model_used": "none",
                        "gpt41_model": query_analysis.model_name,
                        "reasoning": "No documents available"
                    },
                    "query_type": "no_documents",
                    "confidence_score": 0.0,
                    "sources_used": [],
                    "conversation_context": True,
                    "session_id": session_id
                }
            
            # GPT-4.1 ENHANCED: Document search with larger context capacity
            try:
                # ENHANCED: More search results for GPT-4.1's massive context
                search_k = min({1: 8, 2: 12, 3: 20, 4: 35, 5: 50}[query_analysis.level], 80)
                
                search_results = self.vectorstore.similarity_search_with_score(query, k=search_k * 2)
                
                filtered_docs = [(doc, score) for doc, score in search_results 
                               if doc.metadata.get('source', '') != 'init']
                
                if not filtered_docs:
                    return self._generate_no_context_response(query, query_analysis, session_id)
                
                documents = [doc for doc, score in filtered_docs]
                enhanced_results = self._intelligent_context_relevance_scoring(query, documents, query_analysis)
                
                relevant_results = []
                
                for result in enhanced_results:
                    if result.final_score >= self.relevance_thresholds["minimal_relevance"]:
                        relevant_results.append(result)
                    if len(relevant_results) >= search_k:
                        break
                
                if not relevant_results:
                    return self._generate_no_context_response(query, query_analysis, session_id)
                
            except Exception as e:
                return self._generate_error_response(query, query_analysis, session_id, str(e))
            
            # GPT-4.1 OPTIMIZED: Prepare massive context for 1M token capability
            context_parts = []
            sources_used = set()
            total_context_length = 0
            max_context_length = self.context_length_configs[query_analysis.level]
            
            for result in relevant_results:
                doc = result.document
                source = doc.metadata.get('source', 'Unknown')
                content = doc.page_content
                
                if total_context_length + len(content) <= max_context_length:
                    if result.final_score >= self.relevance_thresholds["high_relevance"]:
                        context_parts.append(f"[HIGH RELEVANCE - {source}]\n{content}")
                    elif result.final_score >= self.relevance_thresholds["medium_relevance"]:
                        context_parts.append(f"[RELEVANT - {source}]\n{content}")
                    else:
                        context_parts.append(f"[CONTEXT - {source}]\n{content}")
                    
                    sources_used.add(source)
                    total_context_length += len(content)
                else:
                    break
            
            if not context_parts:
                return self._generate_no_context_response(query, query_analysis, session_id)
            
            context = "\n\n".join(context_parts)
            
            # GPT-4.1 OPTIMIZED: Create enhanced prompt
            enhanced_prompt = self._create_enhanced_prompt(query, context, query_analysis, conversation_context)
            
            # GPT-4.1 OPTIMIZED: Generate response
            try:
                model_config = self.model_configs[query_analysis.level]
                
                llm = ChatOpenAI(
                    model=model_config["model_name"],
                    temperature=model_config["temperature"],
                    max_tokens=model_config["max_tokens"],
                    openai_api_key=self.openai_api_key
                )
                
                response_result = llm.invoke(enhanced_prompt)
                response = response_result.content
                
                # GPT-4.1 OPTIMIZED: Adaptive response enhancement
                response = self._adaptive_response_enhancement(query, response, query_analysis, context, llm)
                
            except Exception as e:
                return self._generate_error_response(query, query_analysis, session_id, str(e))
            
            response_time = time.time() - start_time
            
            # GPT-4.1: Calculate costs and update statistics
            input_tokens = len(self.token_encoder.encode(enhanced_prompt))
            output_tokens = len(self.token_encoder.encode(response))
            
            model_config = self.model_configs[query_analysis.level]
            actual_cost = (
                (input_tokens / 1_000_000) * model_config["input_cost_per_1m"] +
                (output_tokens / 1_000_000) * model_config["output_cost_per_1m"]
            )
            
            # GPT-4.1 OPTIMIZED: Update routing statistics
            self.routing_stats["total_queries"] += 1
            self.routing_stats["level_distribution"][query_analysis.level] += 1
            self.routing_stats["total_cost"] += actual_cost
            self.routing_stats["total_tokens"]["input"] += input_tokens
            self.routing_stats["total_tokens"]["output"] += output_tokens
            
            # NEW: Track GPT-4.1 model usage
            if query_analysis.model_name == "gpt-4.1-nano":
                self.routing_stats["gpt41_nano_queries"] += 1
            elif query_analysis.model_name == "gpt-4.1-mini":
                self.routing_stats["gpt41_mini_queries"] += 1
            elif query_analysis.model_name == "gpt-4.1":
                self.routing_stats["gpt41_full_queries"] += 1
            
            if query_analysis.explicit_detail_request:
                self.routing_stats["explicit_detail_requests"] += 1
            
            if query_analysis.query_type == "follow_up":
                self.routing_stats["follow_up_queries"] += 1
            
            # Calculate confidence score
            avg_relevance = sum(r.final_score for r in relevant_results) / len(relevant_results)
            confidence_score = min(avg_relevance * 1.2, 0.95)
            
            # Generate follow-up suggestions
            follow_up_suggestions = self._generate_follow_up_suggestions(query, query_analysis, list(sources_used))
            
            # GPT-4.1 OPTIMIZED: Update conversation context
            self._update_conversation_context(session_id, query, response, query_analysis, list(sources_used))
            
            # GPT-4.1: Prepare comprehensive response
            enhanced_response = {
                "response": response,
                "complexity_analysis": {
                    "level": query_analysis.level,
                    "level_name": query_analysis.level_name,
                    "model_used": model_config["model_name"],
                    "gpt41_series": True,
                    "reasoning": query_analysis.reasoning,
                    "estimated_cost": query_analysis.estimated_cost,
                    "actual_cost": actual_cost,
                    "response_time": response_time,
                    "context_aware": query_analysis.context_aware,
                    "requires_long_response": query_analysis.requires_long_response,
                    "force_detailed_response": query_analysis.force_detailed_response,
                    "explicit_detail_request": query_analysis.explicit_detail_request,
                    "max_context_tokens": model_config["max_context_tokens"]
                },
                "confidence_score": confidence_score,
                "sources_used": list(sources_used),
                "key_concepts": query_analysis.key_concepts,
                "query_type": query_analysis.query_type,
                "academic_type": query_analysis.academic_type,
                "search_results_count": len(relevant_results),
                "conversation_context": {
                    "session_id": session_id,
                    "total_exchanges": conversation_context.total_exchanges + 1,
                    "context_used": query_analysis.context_aware,
                    "immediate_context_applied": len(conversation_context.conversation_history) > 0
                },
                "follow_up_suggestions": follow_up_suggestions,
                "tokens_used": {
                    "input": input_tokens,
                    "output": output_tokens,
                    "total": input_tokens + output_tokens
                },
                "relevance_scores": {
                    "best_match": relevant_results[0].final_score,
                    "average_relevance": avg_relevance,
                    "total_chunks_analyzed": len(enhanced_results)
                },
                "response_quality": {
                    "length_characters": len(response),
                    "length_words": len(response.split()),
                    "complexity_matched": query_analysis.level,
                    "optimized_length": True,
                    "explicit_detail_handled": query_analysis.explicit_detail_request,
                    "gpt41_optimized": True
                },
                "gpt41_features": {
                    "model_series": "GPT-4.1",
                    "context_capacity_used": len(context),
                    "max_context_capacity": max_context_length,
                    "massive_context_leveraged": len(context) > 10000
                }
            }
            
            return enhanced_response
            
        except Exception as e:
            return self._generate_error_response(query, query_analysis if 'query_analysis' in locals() else None, 
                                               session_id, str(e))

    def _generate_no_context_response(self, query: str, query_analysis: QueryComplexityAnalysis, session_id: str) -> Dict[str, Any]:
        """Generate appropriate response when no relevant context is found"""
        
        if query_analysis.query_type == "follow_up":
            response = "I'd be happy to provide more information, but I need to find relevant content in the uploaded documents first. Could you try rephrasing your question or asking about a different topic that might be covered in your materials?"
        elif query_analysis.query_type == "definition":
            response = f"I don't have a definition for the terms in your query in the uploaded documents. Please check if the topic is covered in your materials, or try asking about related concepts that might be included."
        elif query_analysis.academic_type == "question_paper":
            response = "I don't find any exam questions or question papers in the uploaded documents that match your query. Please make sure you've uploaded the relevant academic materials."
        else:
            response = "I cannot find information relevant to your query in the uploaded documents. This could mean:\n\n• The topic isn't covered in your uploaded materials\n• You might need to rephrase your question\n• Try asking about topics that are definitely in your documents\n\nWould you like to try a different question? 🤔"
        
        self._update_conversation_context(session_id, query, response, query_analysis, [])
        
        return {
            "response": response,
            "complexity_analysis": {
                "level": query_analysis.level,
                "level_name": query_analysis.level_name,
                "model_used": "contextual_response",
                "gpt41_model": query_analysis.model_name,
                "reasoning": "No relevant context found"
            },
            "query_type": query_analysis.query_type,
            "confidence_score": 0.1,
            "sources_used": [],
            "conversation_context": {"session_id": session_id},
            "follow_up_suggestions": [
                "Try asking about a topic you know is in your documents",
                "Rephrase your question with different keywords",
                "Ask about general concepts related to your materials"
            ]
        }

    def _generate_error_response(self, query: str, query_analysis: Optional[QueryComplexityAnalysis], 
                                session_id: str, error_msg: str) -> Dict[str, Any]:
        """Generate appropriate error response"""
        response = f"I encountered an issue while processing your request: {error_msg}. Please try again, and if the problem persists, try rephrasing your question."
        
        if query_analysis:
            self._update_conversation_context(session_id, query, response, query_analysis, [])
        
        return {
            "response": response,
            "complexity_analysis": {
                "level": query_analysis.level if query_analysis else 1,
                "level_name": query_analysis.level_name if query_analysis else "Error",
                "model_used": "error_handler",
                "gpt41_model": query_analysis.model_name if query_analysis else "gpt-4.1-nano",
                "reasoning": f"Error occurred: {error_msg}"
            },
            "query_type": "error",
            "confidence_score": 0.0,
            "sources_used": [],
            "conversation_context": {"session_id": session_id}
        }

    def _generate_follow_up_suggestions(self, query: str, query_analysis: QueryComplexityAnalysis, sources: List[str]) -> List[str]:
        """Generate intelligent follow-up suggestions"""
        suggestions = []
        
        if query_analysis.query_type == "definition":
            suggestions.extend([
                "Can you provide examples of this concept?",
                "How does this relate to other topics?",
                "What are the practical applications?"
            ])
        elif query_analysis.query_type == "explanation":
            suggestions.extend([
                "Can you elaborate on this topic?",
                "What are some examples?",
                "Are there any related concepts I should know?"
            ])
        elif query_analysis.query_type == "comparison":
            suggestions.extend([
                "Which approach is better and why?",
                "What are the advantages and disadvantages?",
                "Can you provide more examples of each?"
            ])
        
        if query_analysis.level <= 2:
            suggestions.extend([
                "Can you explain this in more detail?",
                "What else should I know about this topic?"
            ])
        elif query_analysis.level >= 4:
            suggestions.extend([
                "Are there any related advanced topics?",
                "What are the latest developments in this area?"
            ])
        
        if sources:
            suggestions.append(f"What other topics are covered in {sources[0]}?")
        
        return suggestions[:4]
    


    def generate_group_response(self, query: str, group_context: Dict[str, Any], session_id: str = "default") -> Dict[str, Any]:
        """
        GPT-4.1 OPTIMIZED: Generate response specifically for group study discussions
        Maintains separate context from individual chats
        """
        start_time = time.time()
        
        try:
            # Create group-specific session ID
            group_session_id = f"group_{group_context['group_id']}_thread_{group_context['thread_id']}"
            
            # Get group conversation context
            group_conversation_context = self._get_or_create_group_context(group_session_id, group_context)
            
            # GPT-4.1 OPTIMIZED: Enhanced complexity analysis for group discussions
            query_analysis = self._analyze_group_query_complexity(query, group_conversation_context, group_context)
            
            # FIXED: Use the SAME vectorstore check as normal chat
            if not self.vectorstore:
                # Try to reinitialize if vectorstore is None
                self._initialize_vectorstore()
            
            # Enhanced check - same logic as normal chat should use
            try:
                if not self.vectorstore:
                    return self._generate_group_no_documents_response(group_session_id, group_context, query_analysis, query)
                    
                # Check if vectorstore has actual content (not just init document)
                total_docs = getattr(self.vectorstore.index, 'ntotal', 0) if hasattr(self.vectorstore, 'index') else 0
                
                # More flexible check - if we have processed files, vectorstore should work
                processed_files = self.get_processed_files()
                
                if total_docs <= 1 and not processed_files:
                    return self._generate_group_no_documents_response(group_session_id, group_context, query_analysis, query)
                    
                # If we have processed files but low document count, still try to search
                if not processed_files:
                    return self._generate_group_no_documents_response(group_session_id, group_context, query_analysis, query)
                    
            except Exception as vectorstore_error:
                return self._generate_group_no_documents_response(group_session_id, group_context, query_analysis, query)
            
            # GPT-4.1 ENHANCED: Document search (same as individual but with group weighting)
            try:
                # Enhanced search for group discussions - SAME AS NORMAL CHAT
                search_k = min({1: 10, 2: 15, 3: 25, 4: 40, 5: 60}[query_analysis.level], 100)
                
                # Use the EXACT same search method as normal chat
                search_results = self.vectorstore.similarity_search_with_score(query, k=search_k * 2)
                
                # Apply the SAME filtering as normal chat
                filtered_docs = [(doc, score) for doc, score in search_results 
                            if doc.metadata.get('source', '') != 'init']
                
                if not filtered_docs:
                    return self._generate_group_no_context_response(query, query_analysis, group_session_id, group_context)
                
                documents = [doc for doc, score in filtered_docs]
                enhanced_results = self._intelligent_group_context_scoring(query, documents, query_analysis, group_context)
                
                relevant_results = []
                
                for result in enhanced_results:
                    if result.final_score >= self.relevance_thresholds["minimal_relevance"]:
                        relevant_results.append(result)
                    if len(relevant_results) >= search_k:
                        break
                
                if not relevant_results:
                    return self._generate_group_no_context_response(query, query_analysis, group_session_id, group_context)
                
            except Exception as search_error:
                return self._generate_group_error_response(query, query_analysis, group_session_id, group_context, f"Document search failed: {str(search_error)}")
            
            # GPT-4.1 OPTIMIZED: Prepare context for group discussion
            context_parts = []
            sources_used = set()
            total_context_length = 0
            max_context_length = self.context_length_configs[query_analysis.level]
            
            for result in relevant_results:
                doc = result.document
                source = doc.metadata.get('source', 'Unknown')
                content = doc.page_content
                
                if total_context_length + len(content) <= max_context_length:
                    if result.final_score >= self.relevance_thresholds["high_relevance"]:
                        context_parts.append(f"[HIGH RELEVANCE - {source}]\n{content}")
                    elif result.final_score >= self.relevance_thresholds["medium_relevance"]:
                        context_parts.append(f"[RELEVANT - {source}]\n{content}")
                    else:
                        context_parts.append(f"[CONTEXT - {source}]\n{content}")
                    
                    sources_used.add(source)
                    total_context_length += len(content)
                else:
                    break
            
            if not context_parts:
                return self._generate_group_no_context_response(query, query_analysis, group_session_id, group_context)
            
            context = "\n\n".join(context_parts)
            
            # GPT-4.1 OPTIMIZED: Create GROUP-SPECIFIC enhanced prompt
            enhanced_prompt = self._create_group_enhanced_prompt(query, context, query_analysis, group_conversation_context, group_context)
            
            # GPT-4.1 OPTIMIZED: Generate group response
            try:
                model_config = self.model_configs[query_analysis.level]
                
                llm = ChatOpenAI(
                    model=model_config["model_name"],
                    temperature=model_config["temperature"] + 0.05,  # Slightly more creative for group discussions
                    max_tokens=model_config["max_tokens"],
                    openai_api_key=self.openai_api_key
                )
                
                response_result = llm.invoke(enhanced_prompt)
                response = response_result.content
                
                # GPT-4.1 OPTIMIZED: Group-specific response enhancement
                response = self._adaptive_group_response_enhancement(query, response, query_analysis, context, llm, group_context)
                
            except Exception as llm_error:
                return self._generate_group_error_response(query, query_analysis, group_session_id, group_context, f"Response generation failed: {str(llm_error)}")
            
            response_time = time.time() - start_time
            
            # GPT-4.1: Calculate costs and update statistics
            input_tokens = len(self.token_encoder.encode(enhanced_prompt))
            output_tokens = len(self.token_encoder.encode(response))
            
            model_config = self.model_configs[query_analysis.level]
            actual_cost = (
                (input_tokens / 1_000_000) * model_config["input_cost_per_1m"] +
                (output_tokens / 1_000_000) * model_config["output_cost_per_1m"]
            )
            
            # Update group-specific routing statistics
            self.routing_stats["total_queries"] += 1
            self.routing_stats["level_distribution"][query_analysis.level] += 1
            self.routing_stats["total_cost"] += actual_cost
            
            # Calculate confidence score
            avg_relevance = sum(r.final_score for r in relevant_results) / len(relevant_results)
            confidence_score = min(avg_relevance * 1.2, 0.95)
            
            # Generate group-specific follow-up suggestions
            follow_up_suggestions = self._generate_group_follow_up_suggestions(query, query_analysis, list(sources_used), group_context)
            
            # Update group conversation context
            self._update_group_conversation_context(group_session_id, query, response, query_analysis, list(sources_used), group_context)
            
            # GPT-4.1: Prepare comprehensive GROUP response
            enhanced_response = {
                "response": response,
                "complexity_analysis": {
                    "level": query_analysis.level,
                    "level_name": query_analysis.level_name,
                    "model_used": model_config["model_name"],
                    "gpt41_series": True,
                    "reasoning": query_analysis.reasoning,
                    "estimated_cost": query_analysis.estimated_cost,
                    "actual_cost": actual_cost,
                    "response_time": response_time,
                    "group_optimized": True
                },
                "confidence_score": confidence_score,
                "sources_used": list(sources_used),
                "key_concepts": query_analysis.key_concepts,
                "query_type": query_analysis.query_type,
                "academic_type": query_analysis.academic_type,
                "search_results_count": len(relevant_results),
                "group_context": {
                    "group_id": group_context['group_id'],
                    "group_name": group_context['group_name'],
                    "thread_id": group_context['thread_id'],
                    "thread_title": group_context['thread_title'],
                    "member_count": group_context['member_count'],
                    "session_id": group_session_id,
                    "total_exchanges": group_conversation_context.total_exchanges + 1,
                    "is_group_discussion": True
                },
                "follow_up_suggestions": follow_up_suggestions,
                "tokens_used": {
                    "input": input_tokens,
                    "output": output_tokens,
                    "total": input_tokens + output_tokens
                },
                "relevance_scores": {
                    "best_match": relevant_results[0].final_score,
                    "average_relevance": avg_relevance,
                    "total_chunks_analyzed": len(enhanced_results)
                },
                "response_quality": {
                    "length_characters": len(response),
                    "length_words": len(response.split()),
                    "complexity_matched": query_analysis.level,
                    "group_discussion_optimized": True,
                    "gpt41_optimized": True
                },
                "is_group_response": True,
                "gpt41_features": {
                    "model_series": "GPT-4.1",
                    "context_capacity_used": len(context),
                    "max_context_capacity": max_context_length,
                    "group_discussion_mode": True
                },
                "debug_info": {
                    "vectorstore_exists": bool(self.vectorstore),
                    "processed_files_count": len(processed_files),
                    "total_documents": total_docs,
                    "search_successful": True
                }
            }
            
            return enhanced_response
            
        except Exception as e:
            return self._generate_group_error_response(query, query_analysis if 'query_analysis' in locals() else None, 
                                                    group_session_id if 'group_session_id' in locals() else "error", 
                                                    group_context, str(e))

    def _get_or_create_group_context(self, group_session_id: str, group_context: Dict[str, Any]) -> GroupConversationContext:
        """Get or create group conversation context"""
        with self.group_conversation_lock:
            if group_session_id not in self.group_conversation_contexts:
                self.group_conversation_contexts[group_session_id] = GroupConversationContext(
                    group_id=group_context['group_id'],
                    group_name=group_context['group_name'],
                    thread_id=group_context['thread_id'],
                    thread_title=group_context['thread_title'],
                    member_count=group_context['member_count'],
                    conversation_history=[],
                    last_query="",
                    last_response="",
                    context_documents=[],
                    query_intent="",
                    group_context=group_context,
                    timestamp=datetime.now(),
                    total_exchanges=0,
                    group_session_id=group_session_id
                )
            return self.group_conversation_contexts[group_session_id]

    def _analyze_group_query_complexity(self, query: str, group_context: GroupConversationContext, 
                                    context_info: Dict[str, Any]) -> QueryComplexityAnalysis:
        """
        ENHANCED: Query complexity analysis optimized for GROUP discussions
        """
        # Use the existing complexity analysis but with group-specific adjustments
        base_analysis = self._analyze_advanced_query_complexity(query, None)  # No individual context
        
        # Adjust for group context - group discussions tend to be more collaborative
        if base_analysis.level == 1 and any(word in query.lower() for word in ['explain', 'help', 'understand']):
            base_analysis.level = 2  # Slightly more detailed for group learning
        
        # Group discussions benefit from examples and collaborative explanations
        if any(word in query.lower() for word in ['example', 'show us', 'help everyone', 'group']):
            base_analysis.level = min(base_analysis.level + 1, 5)
        
        return base_analysis

    def _create_group_enhanced_prompt(self, query: str, context: str, query_analysis: QueryComplexityAnalysis,
                                    group_conversation_context: GroupConversationContext, group_context: Dict[str, Any]) -> str:
        """ENHANCED: Create GROUP-SPECIFIC prompts optimized for GPT-4.1 series"""
        
        # Enhanced system prompt for GROUP discussions
        system_prompt = f"""You are SANA, an advanced academic AI assistant powered by {query_analysis.model_name} specifically designed for GROUP STUDY DISCUSSIONS. You're helping a study group collaborate and learn together.

    GROUP STUDY CONTEXT:
    - Group: {group_context['group_name']} (Thread: {group_context['thread_title']})
    - Members: {group_context['member_count']} students
    - Model: {query_analysis.model_name}
    - Query Complexity Level: {query_analysis.level} ({query_analysis.level_name})
    - Group Discussion Mode: ACTIVE
    - Context: Group Study & Collaborative Learning

    GPT-4.1 GROUP DISCUSSION OPTIMIZATION:
    - Encourage collaborative thinking and discussion
    - Provide explanations that benefit the entire group
    - Use inclusive language ("you all", "the group", "everyone")
    - Suggest group activities and discussions when appropriate
    - Be encouraging and supportive of group learning
    """

        # Add group conversation history
        if group_conversation_context.conversation_history:
            recent_history = group_conversation_context.conversation_history[-3:]
            system_prompt += f"""
    GROUP CONVERSATION HISTORY:
    Total group exchanges: {len(group_conversation_context.conversation_history)}
    """
            if recent_history:
                system_prompt += "Recent Group Discussion:\n"
                for i, exchange in enumerate(recent_history, 1):
                    system_prompt += f"Group Q{i}: {exchange['query'][:100]}...\nGroup A{i}: {exchange['response'][:150]}...\n"

        # Group-specific instructions based on complexity level
        group_level_instructions = {
            1: {
                "style": "Provide a clear, direct answer that's perfect for quick group reference. Keep it concise but ensure everyone can understand.",
                "format": "Use simple language perfect for group sharing. Bold key facts with **text**.",
                "group_tone": "Quick and clear for the group! Here's what you need to know:"
            },
            2: {
                "style": "Give a solid explanation that the whole group can follow and discuss together.",
                "format": "Structure your response for easy group discussion. Use bullet points when helpful.",
                "group_tone": "Great question for group discussion! Let me explain this for everyone:"
            },
            3: {
                "style": "Provide a comprehensive explanation perfect for group study sessions with examples the group can work through together.",
                "format": "Use clear sections, examples, and discussion points. Encourage group interaction.",
                "group_tone": "Excellent topic for your study group! Here's a detailed explanation with examples you can explore together:"
            },
            4: {
                "style": "Deliver an in-depth analysis perfect for advanced group study with multiple perspectives and discussion points.",
                "format": "Create structured sections with analysis points, examples, and questions for group discussion.",
                "group_tone": "This is a complex topic perfect for deep group analysis! Let's break this down comprehensively:"
            },
            5: {
                "style": "Provide expert-level comprehensive coverage ideal for intensive group study sessions with research-quality depth.",
                "format": "Structure with clear headings, comprehensive analysis, multiple examples, and extensive discussion opportunities.",
                "group_tone": "Fantastic advanced question for your study group! Here's a comprehensive analysis you can dive deep into together:"
            }
        }
        
        current_instructions = group_level_instructions[query_analysis.level]
        
        # Construct GROUP-OPTIMIZED prompt
        final_prompt = f"""{system_prompt}

    GROUP RESPONSE REQUIREMENTS:
    - Style: {current_instructions['style']}
    - Format: {current_instructions['format']}
    - Tone: {current_instructions['group_tone']}
    - Audience: Study group of {group_context['member_count']} students
    - Context: Collaborative learning environment

    CRITICAL GROUP DISCUSSION RULES:
    1. Answer based ONLY on the provided academic context
    2. Use inclusive language suitable for group learning ("you all", "everyone", "the group")
    3. Encourage group discussion and collaboration when appropriate
    4. Use **bold formatting** for key terms and concepts
    5. Suggest follow-up questions or activities for the group when relevant
    6. Be encouraging and supportive of collaborative learning
    7. If suggesting examples or exercises, make them group-friendly
    8. Structure responses to facilitate group discussion and note-taking

    ACADEMIC CONTEXT FOR GROUP STUDY:
    {context}

    GROUP STUDY QUERY: {query}

    PROVIDE YOUR GROUP-OPTIMIZED ACADEMIC RESPONSE:"""
        
        return final_prompt

    def _adaptive_group_response_enhancement(self, query: str, response: str, query_analysis: QueryComplexityAnalysis,
                                        context: str, llm: ChatOpenAI, group_context: Dict[str, Any]) -> str:
        """Group-specific response enhancement"""
        
        # Check if response needs group-oriented enhancement
        if not any(phrase in response.lower() for phrase in ['group', 'everyone', 'you all', 'together']):
            # Add group-oriented language
            enhancement_prompt = f"""Enhance this response to be more suitable for a study group discussion:

    GROUP CONTEXT: {group_context['group_name']} study group with {group_context['member_count']} members

    CURRENT RESPONSE:
    {response}

    ENHANCED GROUP RESPONSE: Make this more collaborative and group-oriented while maintaining the same information and accuracy. Use inclusive language and encourage group discussion where appropriate."""

            try:
                enhanced_result = llm.invoke(enhancement_prompt)
                if len(enhanced_result.content) > 50:
                    return enhanced_result.content
            except Exception:
                pass
        
        return response

    def _intelligent_group_context_scoring(self, query: str, documents: List[Document], 
                                        query_analysis: QueryComplexityAnalysis, group_context: Dict[str, Any]) -> List[EnhancedSearchResult]:
        """Group-specific context relevance scoring"""
        # Use the existing intelligent scoring but with group-specific weights
        base_results = self._intelligent_context_relevance_scoring(query, documents, query_analysis)
        
        # Adjust scores for group learning context
        for result in base_results:
            content = result.document.page_content.lower()
            
            # Boost content that's good for group discussions
            if any(term in content for term in ['example', 'exercise', 'practice', 'discussion', 'group', 'collaborative']):
                result.final_score *= 1.1
            
            # Boost educational content
            if any(term in content for term in ['explain', 'definition', 'concept', 'theory', 'principle']):
                result.final_score *= 1.05
            
            result.search_method = "group_intelligent_relevance"
        
        base_results.sort(key=lambda x: x.final_score, reverse=True)
        return base_results

    def _update_group_conversation_context(self, group_session_id: str, query: str, response: str,
                                        query_analysis: QueryComplexityAnalysis, sources: List[str], group_context: Dict[str, Any]):
        """Update group conversation context"""
        with self.group_conversation_lock:
            context = self.group_conversation_contexts.get(group_session_id)
            if context:
                context.conversation_history.append({
                    "query": query,
                    "response": response,
                    "timestamp": datetime.now().isoformat(),
                    "complexity_level": query_analysis.level,
                    "sources": sources,
                    "query_type": query_analysis.query_type,
                    "response_length": len(response),
                    "model_used": query_analysis.model_name,
                    "group_context": group_context
                })
                
                if len(context.conversation_history) > self.group_conversation_memory_limit:
                    context.conversation_history = context.conversation_history[-self.group_conversation_memory_limit:]
                
                context.last_query = query
                context.last_response = response
                context.context_documents = sources
                context.query_intent = query_analysis.query_type
                context.timestamp = datetime.now()
                context.total_exchanges += 1

    def _generate_group_follow_up_suggestions(self, query: str, query_analysis: QueryComplexityAnalysis, 
                                            sources: List[str], group_context: Dict[str, Any]) -> List[str]:
        """Generate group-specific follow-up suggestions"""
        suggestions = []
        
        # Group-oriented suggestions
        if query_analysis.query_type == "definition":
            suggestions.extend([
                "Can you show us some examples for the group to discuss?",
                "What are some practice questions we can work on together?",
                "How can our group apply this concept?"
            ])
        elif query_analysis.query_type == "explanation":
            suggestions.extend([
                "Can you give us some exercises to practice as a group?",
                "What related topics should our study group explore?",
                "Are there any group activities to reinforce this concept?"
            ])
        elif query_analysis.query_type == "comparison":
            suggestions.extend([
                "Can you help us create a comparison chart for group study?",
                "What are some good examples for group discussion?",
                "How can we practice identifying these differences together?"
            ])
        
        # Add group-specific suggestions
        suggestions.extend([
            f"What other topics from {sources[0] if sources else 'our materials'} should we cover in our study session?",
            "Can you suggest some group discussion questions on this topic?",
            "What are the key points our group should focus on for exams?"
        ])
        
        return suggestions[:4]

    def _generate_group_no_context_response(self, query: str, query_analysis: QueryComplexityAnalysis, 
                                        group_session_id: str, group_context: Dict[str, Any]) -> Dict[str, Any]:
        """Generate appropriate response when no relevant context is found for group"""
        
        response = f"Hey {group_context['group_name']} study group! 📚 I couldn't find information relevant to '{query}' in your uploaded study materials. This could mean:\n\n• The topic isn't covered in your current materials\n• Try rephrasing your question with different keywords\n• Ask about topics you know are in your documents\n\nWould someone in the group like to try a different question? I'm here to help with your collaborative study session! 🎓"
        
        self._update_group_conversation_context(group_session_id, query, response, query_analysis, [], group_context)
        
        return {
            "response": response,
            "complexity_analysis": {
                "level": query_analysis.level,
                "level_name": query_analysis.level_name,
                "model_used": "group_contextual_response",
                "gpt41_model": query_analysis.model_name,
                "reasoning": "No relevant context found for group discussion"
            },
            "query_type": query_analysis.query_type,
            "confidence_score": 0.1,
            "sources_used": [],
            "group_context": group_context,
            "is_group_response": True,
            "follow_up_suggestions": [
                "Try asking about a topic you know is in your study materials",
                "Rephrase your question with different keywords",
                "Ask about general concepts related to your coursework",
                "Check what materials have been uploaded for group study"
            ]
        }

    def _generate_group_no_documents_response(self, group_session_id: str, group_context: Dict[str, Any], 
                                            query_analysis: QueryComplexityAnalysis, query: str) -> Dict[str, Any]:
        """Generate response when no documents are available for group"""
        
        processed_files = self.get_processed_files()
        
        if not processed_files:
            response = f"Hey {group_context['group_name']} study group! 📚 I don't have any documents to reference for our discussion. Please upload some study materials first, and I'll be happy to help with group questions and discussions! 🎓"
        else:
            response = f"Hey {group_context['group_name']} study group! 📚 I can see there are {len(processed_files)} files uploaded ({', '.join([f['filename'] for f in processed_files[:3]])}{'...' if len(processed_files) > 3 else ''}), but I'm having trouble accessing them right now. Please try your question again, or contact support if this continues! 🎓"
        
        self._update_group_conversation_context(group_session_id, query, response, query_analysis, [], group_context)
        
        return {
            "response": response,
            "complexity_analysis": {
                "level": query_analysis.level,
                "level_name": query_analysis.level_name,
                "model_used": "none",
                "gpt41_model": query_analysis.model_name,
                "reasoning": "No documents available for group discussion"
            },
            "query_type": "group_no_documents",
            "confidence_score": 0.0,
            "sources_used": [],
            "group_context": group_context,
            "is_group_response": True,
            "session_id": group_session_id,
            "processed_files_count": len(processed_files),
            "debug_info": {
                "vectorstore_exists": bool(self.vectorstore),
                "processed_files": len(processed_files),
                "issue": "vectorstore_access_problem" if processed_files else "no_documents_uploaded"
            }
        }

    def _generate_group_error_response(self, query: str, query_analysis: Optional[QueryComplexityAnalysis],
                                    group_session_id: str, group_context: Dict[str, Any], error_msg: str) -> Dict[str, Any]:
        """Generate appropriate error response for group"""
        response = f"Sorry {group_context['group_name']} study group! I encountered an issue while processing your question: {error_msg}. Please try asking your question again, and if the problem continues, try rephrasing it differently. I'm here to help your group study session! 🤔"
        
        if query_analysis:
            self._update_group_conversation_context(group_session_id, query, response, query_analysis, [], group_context)
        
        return {
            "response": response,
            "complexity_analysis": {
                "level": query_analysis.level if query_analysis else 1,
                "level_name": query_analysis.level_name if query_analysis else "Group Error",
                "model_used": "group_error_handler",
                "gpt41_model": query_analysis.model_name if query_analysis else "gpt-4.1-nano",
                "reasoning": f"Group error occurred: {error_msg}"
            },
            "query_type": "group_error",
            "confidence_score": 0.0,
            "sources_used": [],
            "group_context": group_context,
            "is_group_response": True
        }

    def clear_group_conversation_memory(self, group_session_id: str):
        """Clear group conversation memory for session"""
        with self.group_conversation_lock:
            if group_session_id in self.group_conversation_contexts:
                del self.group_conversation_contexts[group_session_id]

    def get_group_conversation_stats(self) -> Dict[str, Any]:
        """Get group conversation statistics"""
        return {
            "active_group_sessions": len(self.group_conversation_contexts),
            "group_conversations": {
                session_id: {
                    "group_id": context.group_id,
                    "group_name": context.group_name,
                    "thread_title": context.thread_title,
                    "total_exchanges": context.total_exchanges,
                    "last_activity": context.timestamp.isoformat()
                }
                for session_id, context in self.group_conversation_contexts.items()
            },
            "group_features_active": True,
            "separate_group_context": True
        }

    def clear_all_group_memories(self):
        """Clear all group conversation memories"""
        with self.group_conversation_lock:
            self.group_conversation_contexts.clear()


    # ALL DOCUMENT PROCESSING METHODS
    
    def _process_pdf_advanced(self, filepath: str) -> str:
        """Advanced PDF processing with enhanced text extraction"""
        all_text = ""
        
        try:
            doc = fitz.open(filepath)
            
            for page_num in range(len(doc)):
                try:
                    page = doc.load_page(page_num)
                    
                    text_dict = page.get_text("dict")
                    page_text = ""
                    
                    for block in text_dict.get("blocks", []):
                        if "lines" in block:
                            for line in block["lines"]:
                                line_text = ""
                                for span in line.get("spans", []):
                                    text = span.get("text", "").strip()
                                    if text:
                                        flags = span.get("flags", 0)
                                        if flags & 2**4:
                                            text = f"**{text}**"
                                        elif flags & 2**1:
                                            text = f"*{text}*"
                                        line_text += text + " "
                                
                                if line_text.strip():
                                    page_text += line_text.strip() + "\n"
                    
                    try:
                        tables = page.find_tables()
                        for table_idx, table in enumerate(tables):
                            table_data = table.extract()
                            if table_data:
                                page_text += f"\n[TABLE {table_idx + 1}]\n"
                                for row in table_data:
                                    clean_row = [str(cell).strip() for cell in row if cell]
                                    if clean_row:
                                        page_text += " | ".join(clean_row) + "\n"
                                page_text += "\n"
                    except Exception:
                        pass
                    
                    if page_text.strip():
                        all_text += f"\n--- Page {page_num + 1} ---\n{page_text}"
                
                except Exception as e:
                    continue
            
            doc.close()
            
        except Exception as e:
            try:
                with pdfplumber.open(filepath) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        if page_text:
                            all_text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
            except Exception as fallback_error:
                pass
        
        return all_text

    def _process_docx_advanced(self, filepath: str) -> str:
        """Advanced DOCX processing with structure preservation"""
        all_text = ""
        
        try:
            doc = docx.Document(filepath)
            
            props = doc.core_properties
            if props.title:
                all_text += f"Title: {props.title}\n\n"
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text = para.text.strip()
                    style = para.style.name
                    
                    if style.startswith('Heading'):
                        level = style.replace('Heading ', '')
                        all_text += f"\n{'#' * int(level if level.isdigit() else '1')} {text}\n\n"
                    else:
                        all_text += f"{text}\n"
            
            for table in doc.tables:
                all_text += "\n[TABLE]\n"
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        all_text += " | ".join(row_text) + "\n"
                all_text += "\n"
            
        except Exception as e:
            pass
        
        return all_text

    def _process_excel_advanced(self, filepath: str) -> str:
        """Advanced Excel processing with enhanced academic content detection"""
        all_text = ""
        try:
            if os.path.basename(filepath).startswith('~$'):
                return ""
            
            excel_file = pd.ExcelFile(filepath, engine="openpyxl")
            sheet_names = excel_file.sheet_names
            
            for sheet_name in sheet_names:
                try:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name, engine="openpyxl")
                    if not df.empty:
                        all_text += f"\n=== EXCEL SHEET: {sheet_name} ===\n"
                        all_text += f"Data Dimensions: {df.shape[0]} rows × {df.shape[1]} columns\n"
                        all_text += f"Columns: {', '.join(df.columns.astype(str))}\n\n"
                        
                        if len(df) > 1000:
                            all_text += "=== DATA SAMPLE (First 500 rows) ===\n"
                            sample_df = df.head(500)
                        else:
                            all_text += "=== COMPLETE DATA ===\n"
                            sample_df = df
                        
                        all_text += sample_df.to_string(index=True, max_rows=None, max_cols=None) + "\n\n"
                        
                        numeric_cols = df.select_dtypes(include=[np.number]).columns
                        if len(numeric_cols) > 0:
                            all_text += f"=== STATISTICS ===\n"
                            all_text += df[numeric_cols].describe().to_string() + "\n\n"
                        
                        for col in df.columns:
                            if any(term in str(col).lower() for term in ['grade', 'marks', 'score', 'result']):
                                all_text += f"Academic Data Column: {col}\n"
                
                except Exception as sheet_error:
                    pass
            
        except Exception as e:
            pass
        
        return all_text

    def _process_powerpoint_advanced(self, filepath: str) -> str:
        """Advanced PowerPoint processing with enhanced academic structure detection"""
        all_text = ""
        try:
            prs = Presentation(filepath)
            
            if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
                all_text += f"=== PRESENTATION TITLE ===\n{prs.core_properties.title}\n\n"
            
            slides = list(prs.slides)
            for slide_num, slide in enumerate(slides):
                all_text += f"\n=== SLIDE {slide_num + 1} ===\n"
                
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        all_text += f"Slide Title: {title_text}\n\n"
                
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape != slide.shapes.title:
                            text = shape.text.strip()
                            
                            if self._is_question_pattern(text):
                                all_text += f"[Question] {text}\n"
                            elif self._is_answer_pattern(text):
                                all_text += f"[Answer] {text}\n"
                            elif any(term in text.lower() for term in ['definition', 'theorem', 'principle']):
                                all_text += f"[Definition] {text}\n"
                            elif any(term in text.lower() for term in ['example', 'case study', 'illustration']):
                                all_text += f"[Example] {text}\n"
                            else:
                                all_text += f"[Content] {text}\n"
                
                if slide.has_notes_slide:
                    try:
                        notes_slide = slide.notes_slide
                        if notes_slide.notes_text_frame:
                            notes_text = notes_slide.notes_text_frame.text.strip()
                            if notes_text:
                                all_text += f"[Notes] {notes_text}\n"
                    except Exception as notes_error:
                        pass
                
                all_text += "\n"
            
        except Exception as e:
            pass
        
        return all_text

    def _process_text_advanced(self, filepath: str) -> str:
        """Advanced text file processing with enhanced academic structure detection"""
        all_text = ""
        try:
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'ascii', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(filepath, 'r', encoding=encoding) as file:
                        content = file.read()
                    
                    if content.strip():
                        lines = content.split('\n')
                        processed_lines = []
                        
                        all_text += f"=== TEXT FILE ANALYSIS ===\n"
                        all_text += f"File: {os.path.basename(filepath)}\n"
                        all_text += f"Size: {len(content)} characters\n"
                        all_text += f"Lines: {len(lines)}\n"
                        all_text += f"Encoding: {encoding}\n\n"
                        
                        question_count = 0
                        answer_count = 0
                        definition_count = 0
                        
                        for line_num, line in enumerate(lines, 1):
                            line = line.strip()
                            if line:
                                if self._is_question_pattern(line):
                                    processed_lines.append(f"[Question - Line {line_num}] {line}")
                                    question_count += 1
                                elif self._is_answer_pattern(line):
                                    processed_lines.append(f"[Answer - Line {line_num}] {line}")
                                    answer_count += 1
                                elif any(term in line.lower() for term in ['definition', 'define', 'means', 'refers to']):
                                    processed_lines.append(f"[Definition - Line {line_num}] {line}")
                                    definition_count += 1
                                elif line.startswith('#') or line.isupper():
                                    processed_lines.append(f"[Header - Line {line_num}] {line}")
                                elif any(char in line for char in ['•', '-', '*']) and len(line.split()) < 20:
                                    processed_lines.append(f"[List Item - Line {line_num}] {line}")
                                else:
                                    processed_lines.append(f"[Content - Line {line_num}] {line}")
                        
                        all_text += f"=== DOCUMENT STRUCTURE ===\n"
                        all_text += f"Questions Found: {question_count}\n"
                        all_text += f"Answers Found: {answer_count}\n"
                        all_text += f"Definitions Found: {definition_count}\n\n"
                        
                        all_text += "=== CONTENT ===\n"
                        all_text += '\n'.join(processed_lines)
                        break
                        
                except (UnicodeDecodeError, UnicodeError):
                    continue
                except Exception as e:
                    continue
            
        except Exception as e:
            pass
        
        return all_text

    def _process_html_advanced(self, filepath: str) -> str:
        """Advanced HTML processing with enhanced academic content extraction"""
        all_text = ""
        try:
            with open(filepath, 'r', encoding='utf-8') as file:
                content = file.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            for script in soup(["script", "style"]):
                script.decompose()
            
            title = soup.find('title')
            if title:
                all_text += f"=== HTML DOCUMENT ===\n"
                all_text += f"Title: {title.get_text()}\n\n"
            
            meta_desc = soup.find('meta', attrs={'name': 'description'})
            if meta_desc:
                all_text += f"Description: {meta_desc.get('content', '')}\n"
            
            meta_keywords = soup.find('meta', attrs={'name': 'keywords'})
            if meta_keywords:
                all_text += f"Keywords: {meta_keywords.get('content', '')}\n\n"
            
            headings = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
            if headings:
                all_text += "=== DOCUMENT STRUCTURE ===\n"
                for heading in headings:
                    level = heading.name.upper()
                    text = heading.get_text().strip()
                    indent = "  " * (int(heading.name[1]) - 1)
                    all_text += f"{indent}{level}: {text}\n"
                all_text += "\n"
            
            all_text += "=== CONTENT ===\n"
            
            paragraphs = soup.find_all('p')
            for para_idx, para in enumerate(paragraphs):
                para_text = para.get_text().strip()
                if para_text:
                    if self._is_question_pattern(para_text):
                        all_text += f"[Question {para_idx + 1}] {para_text}\n\n"
                    elif self._is_answer_pattern(para_text):
                        all_text += f"[Answer {para_idx + 1}] {para_text}\n\n"
                    elif any(term in para_text.lower() for term in ['definition', 'theorem', 'principle']):
                        all_text += f"[Definition {para_idx + 1}] {para_text}\n\n"
                    else:
                        all_text += f"[Paragraph {para_idx + 1}] {para_text}\n\n"
            
            lists = soup.find_all(['ul', 'ol'])
            for list_idx, list_elem in enumerate(lists):
                list_type = "Ordered" if list_elem.name == 'ol' else "Unordered"
                all_text += f"[{list_type} List {list_idx + 1}]\n"
                
                items = list_elem.find_all('li')
                for item_idx, item in enumerate(items):
                    item_text = item.get_text().strip()
                    if item_text:
                        all_text += f"  • {item_text}\n"
                all_text += "\n"
            
            tables = soup.find_all('table')
            for table_idx, table in enumerate(tables):
                all_text += f"=== TABLE {table_idx + 1} ===\n"
                
                headers = table.find_all('th')
                if headers:
                    header_text = [header.get_text().strip() for header in headers]
                    all_text += f"Headers: {' | '.join(header_text)}\n"
                
                rows = table.find_all('tr')
                for row_idx, row in enumerate(rows, 1):
                    cells = row.find_all(['td', 'th'])
                    if cells and not (len(headers) > 0 and row_idx == 1):
                        row_text = [cell.get_text().strip() for cell in cells]
                        all_text += f"Row {row_idx}: {' | '.join(row_text)}\n"
                all_text += "\n"
            
            code_blocks = soup.find_all(['code', 'pre'])
            for code_idx, code in enumerate(code_blocks):
                code_text = code.get_text().strip()
                if code_text:
                    all_text += f"[Code Block {code_idx + 1}]\n{code_text}\n\n"
            
        except Exception as e:
            pass
        
        return all_text

    def _is_question_pattern(self, text: str) -> bool:
        """Enhanced question pattern detection"""
        text_lower = text.lower()
        question_patterns = [
            r'\bq\.?\s*\d+', r'\bquestion\s*\d+', r'\b\d+\s*marks?\b',
            r'\bwhat\s+(?:is|are|does|do)\b', r'\bwhy\s+(?:is|are|does|do)\b', 
            r'\bhow\s+(?:does|do|can|is|are)\b', r'\bwhen\s+(?:is|are|does|do)\b',
            r'\bwhere\s+(?:is|are|does|do)\b', r'\bwhich\s+(?:is|are|does|do)\b',
            r'\bexplain\b', r'\bdiscuss\b', r'\bdefine\b', r'\bcompare\b',
            r'\banalyze\b', r'\bevaluate\b', r'\bdescribe\b'
        ]
        return any(re.search(pattern, text_lower) for pattern in question_patterns) or text.endswith('?')

    def _is_answer_pattern(self, text: str) -> bool:
        """Enhanced answer pattern detection"""
        text_lower = text.lower()
        answer_patterns = [
            r'\banswer\s*:?', r'\bsolution\s*:?', r'\bexplanation\s*:?',
            r'\bsteps?\s*:?', r'\bmethod\s*:?', r'\bapproach\s*:?',
            r'\bresult\s*:?', r'\bconclusion\s*:?', r'\bfinding\s*:?'
        ]
        return any(re.search(pattern, text_lower) for pattern in answer_patterns)

    # ALL CATEGORIZATION AND PROCESSING METHODS
    
    def _robust_subject_categorization(self, text: str, filename: str) -> DocumentCategory:
        """Robust subject and category detection using multiple strategies"""
        
        text_lower = text.lower()
        text_words = safe_word_tokenize(text_lower)
        
        subject_scores = {}
        
        for subject, config in self.subject_domains.items():
            score = 0.0
            
            for keyword in config["keywords"]:
                count = text_lower.count(keyword.lower())
                if count > 0:
                    score += min(count * 2.0, 10.0)
            
            for pattern in config["patterns"]:
                matches = len(re.findall(pattern, text_lower))
                score += matches * 3.0
            
            if len(text_words) > 0:
                score = score / math.log(len(text_words) + 1)
            
            subject_scores[subject] = score
        
        if len(text_words) > 50:
            try:
                subject_vocabularies = {}
                for subject, config in self.subject_domains.items():
                    vocab = config["keywords"] + [term for pattern in config["patterns"] 
                                                 for term in re.findall(r'\w+', pattern)]
                    subject_vocabularies[subject] = set(vocab)
                
                for subject, vocab in subject_vocabularies.items():
                    vocab_matches = sum(1 for word in text_words if word in vocab)
                    if vocab_matches > 0:
                        tf_score = vocab_matches / len(text_words)
                        subject_scores[subject] += tf_score * 100
                        
            except Exception as e:
                pass
        
        structure_indicators = {
            'textbook': ['chapter', 'section', 'exercise', 'summary', 'review'],
            'research_paper': ['abstract', 'introduction', 'methodology', 'results', 'conclusion', 'references'],
            'lecture_notes': ['lecture', 'slide', 'presentation', 'topic', 'agenda'],
            'assignment': ['assignment', 'homework', 'project', 'task', 'due'],
            'exam_material': ['exam', 'test', 'quiz', 'question', 'marks', 'points']
        }
        
        document_type_scores = {}
        for doc_type, indicators in structure_indicators.items():
            score = sum(text_lower.count(indicator) for indicator in indicators)
            document_type_scores[doc_type] = score
        
        sorted_subjects = sorted(subject_scores.items(), key=lambda x: x[1], reverse=True)
        
        if not sorted_subjects or sorted_subjects[0][1] == 0:
            primary_subject = "general_studies"
            secondary_subjects = []
            confidence = 0.3
        else:
            primary_subject = sorted_subjects[0][0]
            threshold = sorted_subjects[0][1] * 0.3
            secondary_subjects = [subject for subject, score in sorted_subjects[1:4] 
                                if score > threshold]
            confidence = min(sorted_subjects[0][1] / 50.0, 0.95)
        
        if document_type_scores:
            best_doc_type = max(document_type_scores, key=document_type_scores.get)
            if document_type_scores[best_doc_type] > 0:
                document_type = best_doc_type
            else:
                document_type = "general_document"
        else:
            document_type = "general_document"
        
        academic_level_indicators = {
            'elementary': ['basic', 'simple', 'introduction', 'beginner'],
            'intermediate': ['intermediate', 'moderate', 'standard'],
            'advanced': ['advanced', 'complex', 'sophisticated', 'expert'],
            'graduate': ['graduate', 'research', 'thesis', 'dissertation', 'advanced research']
        }
        
        level_scores = {}
        for level, indicators in academic_level_indicators.items():
            score = sum(text_lower.count(indicator) for indicator in indicators)
            level_scores[level] = score
        
        if level_scores and max(level_scores.values()) > 0:
            academic_level = max(level_scores, key=level_scores.get)
        else:
            academic_level = "intermediate"
        
        key_topics = self._extract_key_topics_advanced(text, primary_subject)
        
        content_categories = []
        if primary_subject in self.subject_domains:
            content_categories = self.subject_domains[primary_subject]["subcategories"]
        
        return DocumentCategory(
            primary_subject=primary_subject,
            secondary_subjects=secondary_subjects,
            document_type=document_type,
            academic_level=academic_level,
            content_categories=content_categories,
            key_topics=key_topics,
            confidence_score=confidence,
            language="english",
            complexity_level=academic_level
        )

    def _extract_key_topics_advanced(self, text: str, subject: str) -> List[str]:
        """Extract key topics using advanced NLP techniques"""
        try:
            text_lower = text.lower()
            
            if subject in self.subject_domains:
                domain_keywords = self.subject_domains[subject]["keywords"]
                topics = []
                
                for keyword in domain_keywords:
                    if keyword in text_lower:
                        context_matches = re.findall(
                            r'\b\w*' + re.escape(keyword) + r'\w*\b',
                            text_lower
                        )
                        topics.extend(context_matches[:3])
                
                return list(set(topics))[:10]
            
            words = safe_word_tokenize(text_lower)
            topic_words = [word for word in words 
                          if len(word) > 4 and word not in STOP_WORDS and word.isalpha()]
            
            word_counts = Counter(topic_words)
            return [word for word, count in word_counts.most_common(10)]
            
        except Exception as e:
            return []

    def _advanced_chunking_with_categorization(self, text: str, filename: str, 
                                             document_category: DocumentCategory) -> List[Document]:
        """ENHANCED: Advanced chunking optimized for GPT-4.1's context capabilities"""
        
        chunks = []
        
        try:
            text_length = len(text)
            
            # ENHANCED: Larger chunks for GPT-4.1's massive context windows
            if text_length < 5000:
                chunk_config = {'chunk_size': 500, 'chunk_overlap': 75}
            elif text_length < 20000:
                chunk_config = {'chunk_size': 800, 'chunk_overlap': 150}
            elif text_length < 100000:
                chunk_config = {'chunk_size': 1200, 'chunk_overlap': 250}
            else:
                chunk_config = {'chunk_size': 1500, 'chunk_overlap': 300}
            
            # Adjust for document type
            if document_category.document_type == 'textbook':
                chunk_config['chunk_size'] = int(chunk_config['chunk_size'] * 1.3)
            elif document_category.document_type == 'exam_material':
                chunk_config['chunk_size'] = int(chunk_config['chunk_size'] * 0.9)
            
            # Create enhanced text splitter
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_config['chunk_size'],
                chunk_overlap=chunk_config['chunk_overlap'],
                separators=[
                    "\n\n=== ", "\n\n", "\n", ". ", "! ", "? ", 
                    "; ", ": ", ", ", " ", ""
                ]
            )
            
            base_chunks = splitter.split_text(text)
            
            for i, chunk_content in enumerate(base_chunks):
                if len(chunk_content.strip()) > 50:
                    
                    chunk_academic_type = self._detect_chunk_academic_type(chunk_content)
                    
                    importance_score = self._calculate_advanced_chunk_importance(
                        chunk_content, document_category
                    )
                    
                    chunk = Document(
                        page_content=chunk_content,
                        metadata={
                            'source': filename,
                            'chunk_id': f"{filename}_{i}",
                            'chunk_index': i,
                            'chunk_type': 'gpt41_enhanced',
                            'document_type': document_category.document_type,
                            'primary_subject': document_category.primary_subject,
                            'secondary_subjects': document_category.secondary_subjects,
                            'academic_level': document_category.academic_level,
                            'content_categories': document_category.content_categories,
                            'chunk_academic_type': chunk_academic_type,
                            'importance_score': importance_score,
                            'complexity_level': document_category.complexity_level,
                            'key_topics': document_category.key_topics,
                            'processing_timestamp': datetime.now().isoformat(),
                            'chunk_length': len(chunk_content),
                            'category': asdict(document_category),
                            'gpt41_optimized': True
                        }
                    )
                    chunks.append(chunk)
            
            return chunks
            
        except Exception as e:
            return self._simple_fallback_chunking(text, filename)

    def _calculate_advanced_chunk_importance(self, content: str, 
                                           document_category: DocumentCategory) -> float:
        """Calculate chunk importance with advanced categorization context"""
        
        importance_score = 0.5
        content_lower = content.lower()
        
        academic_indicators = [
            'definition', 'theorem', 'principle', 'concept', 'theory',
            'example', 'formula', 'equation', 'method', 'procedure',
            'important', 'key', 'main', 'primary', 'essential', 'critical'
        ]
        
        for indicator in academic_indicators:
            count = content_lower.count(indicator)
            importance_score += min(count * 0.05, 0.2)
        
        if document_category.primary_subject in self.subject_domains:
            subject_keywords = self.subject_domains[document_category.primary_subject]['keywords']
            for keyword in subject_keywords:
                if keyword.lower() in content_lower:
                    importance_score += 0.03
        
        type_multipliers = {
            'textbook': 1.2,
            'research_paper': 1.3,
            'lecture_notes': 1.1,
            'exam_material': 1.4,
            'assignment': 1.2
        }
        
        multiplier = type_multipliers.get(document_category.document_type, 1.0)
        importance_score *= multiplier
        
        structure_indicators = [
            'chapter', 'section', 'question', 'answer', 'solution',
            'example', 'case study', 'exercise', 'problem'
        ]
        
        structure_count = sum(1 for indicator in structure_indicators 
                            if indicator in content_lower)
        importance_score += min(structure_count * 0.05, 0.15)
        
        if re.search(r'\d+|\$|\%|=|\+|\-|\*|\/|\^', content):
            importance_score += 0.1
        
        return min(importance_score, 1.0)

    def _detect_chunk_academic_type(self, content: str) -> str:
        """Detect specific academic type of chunk content"""
        content_lower = content.lower()
        
        if any(re.search(pattern, content_lower) for pattern in [
            r'\bq\.?\s*\d+', r'\bquestion\s*\d+', r'\b\d+\s*marks?\b'
        ]):
            return 'question'
        
        if any(word in content_lower for word in [
            'answer:', 'solution:', 'explanation:', 'steps:'
        ]):
            return 'answer'
        
        if any(phrase in content_lower for phrase in [
            'is defined as', 'refers to', 'definition', 'meaning'
        ]):
            return 'definition'
        
        if any(word in content_lower for word in [
            'example', 'for instance', 'such as', 'case study'
        ]):
            return 'example'
        
        if any(word in content_lower for word in [
            'formula', 'equation', 'algorithm', 'method', 'procedure'
        ]):
            return 'formula'
        
        return 'content'

    def _detect_academic_content(self, text: str) -> bool:
        """Enhanced academic content detection"""
        text_lower = text.lower()
        
        academic_indicators = {
            'strong': ['academic', 'university', 'college', 'course', 'curriculum', 'syllabus'],
            'moderate': ['study', 'learn', 'education', 'knowledge', 'research', 'analysis'],
            'weak': ['book', 'chapter', 'section', 'page', 'content', 'material']
        }
        
        score = 0
        for strength, indicators in academic_indicators.items():
            for indicator in indicators:
                count = text_lower.count(indicator)
                if strength == 'strong':
                    score += count * 3
                elif strength == 'moderate':
                    score += count * 2
                else:
                    score += count
        
        academic_patterns = [
            r'\bchapter\s+\d+\b', r'\bsection\s+\d+\b', r'\bexercise\s+\d+\b',
            r'\bquestion\s+\d+\b', r'\bexam\b', r'\btest\b', r'\bassignment\b'
        ]
        
        pattern_matches = sum(1 for pattern in academic_patterns 
                            if re.search(pattern, text_lower))
        score += pattern_matches * 5
        
        return score >= 10

    def _simple_fallback_chunking(self, text: str, filename: str) -> List[Document]:
        """Simple fallback chunking method"""
        try:
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
            
            chunks = []
            texts = splitter.split_text(text)
            
            for i, chunk_text in enumerate(texts):
                if len(chunk_text.strip()) > 50:
                    chunk = Document(
                        page_content=chunk_text,
                        metadata={
                            'source': filename,
                            'chunk_id': f"{filename}_fallback_{i}",
                            'chunk_type': 'gpt41_fallback',
                            'importance_score': 0.5,
                            'gpt41_optimized': True
                        }
                    )
                    chunks.append(chunk)
            
            return chunks
            
        except Exception as e:
            return []

    # ALL WORKFLOW AND SYSTEM METHODS
    
    def process_documents(self, file_paths: List[str], progress_callback=None) -> Dict[str, Any]:
        """GPT-4.1 ENHANCED: Document processing with advanced categorization"""
        start_time = time.time()
        
        results = {
            'total_files': len(file_paths),
            'processed_files': 0,
            'successful_files': [],
            'failed_files': [],
            'processing_results': [],
            'total_chunks_added': 0,
            'processing_time': 0.0,
            'academic_content_detected': 0,
            'categories_detected': {},
            'subjects_identified': set(),
            'gpt41_optimized': True
        }
        
        for i, file_path in enumerate(file_paths):
            try:
                filename = os.path.basename(file_path)
                
                if progress_callback:
                    progress_callback({
                        'current_file': filename,
                        'processed': i,
                        'total': len(file_paths),
                        'percentage': (i / len(file_paths)) * 100,
                        'gpt41_processing': True
                    })
                
                result = self.process_document(file_path, filename)
                results['processing_results'].append(result)
                
                if result.success:
                    results['successful_files'].append(filename)
                    results['total_chunks_added'] += result.chunks_added
                    if result.academic_content_detected:
                        results['academic_content_detected'] += 1
                    
                    if result.metadata.get('category'):
                        category_info = result.metadata['category']
                        primary_subject = category_info.get('primary_subject')
                        if primary_subject:
                            results['subjects_identified'].add(primary_subject)
                            if primary_subject not in results['categories_detected']:
                                results['categories_detected'][primary_subject] = []
                            results['categories_detected'][primary_subject].append(filename)
                    
                else:
                    results['failed_files'].append({
                        'filename': filename,
                        'error': result.error
                    })
                
            except Exception as e:
                error_msg = f"Unexpected error processing {filename}: {str(e)}"
                results['failed_files'].append({
                    'filename': filename,
                    'error': error_msg
                })
        
        results['processed_files'] = len(results['successful_files'])
        results['processing_time'] = time.time() - start_time
        results['subjects_identified'] = list(results['subjects_identified'])
        
        try:
            self._update_system_metadata(results)
            self._build_enhanced_search_indices()
        except Exception as e:
            pass
        
        return results

    def process_document(self, filepath: str, filename: str = None) -> ProcessingResult:
        """GPT-4.1 ENHANCED: Single document processing with large document support"""
        start_time = time.time()
        
        if filename is None:
            filename = os.path.basename(filepath)
        
        try:
            
            if not os.path.exists(filepath):
                raise FileNotFoundError(f"File not found: {filepath}")
            
            file_ext = os.path.splitext(filename)[1].lower()
            if file_ext not in self.document_processors:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            processor_func = self.document_processors[file_ext]
            extracted_text = processor_func(filepath)
            
            if not extracted_text or len(extracted_text.strip()) < 20:
                raise ValueError(f"Insufficient content extracted from {filename}")
            
            document_category = self._robust_subject_categorization(extracted_text, filename)
            
            academic_content_detected = self._detect_academic_content(extracted_text)
            
            chunks = self._advanced_chunking_with_categorization(
                extracted_text, filename, document_category
            )
            
            if not chunks:
                raise ValueError(f"No chunks created from {filename}")
            
            if filename in self.document_metadata.get('processed_files', []):
                self.remove_document(filename)
            
            embedding_success = self._process_large_document_embeddings(chunks, filename)
            
            if not embedding_success:
                raise ValueError(f"Failed to process embeddings for {filename}")
            
            processing_time = time.time() - start_time
            
            # GPT-4.1 ENHANCED: Metadata
            file_metadata = {
                'filename': filename,
                'file_type': file_ext,
                'processed_at': datetime.now().isoformat(),
                'chunk_count': len(chunks),
                'content_length': len(extracted_text),
                'processing_time': processing_time,
                'academic_content_detected': academic_content_detected,
                'category': asdict(document_category),
                'embedding_model': 'text-embedding-3-large',
                'embedding_dimensions': 3072,
                'enhanced_processing': True,
                'large_document_processing': len(extracted_text) > self.large_doc_threshold,
                'file_size_bytes': os.path.getsize(filepath),
                'complexity_level': document_category.complexity_level,
                'optimized_response_system': True,
                'gpt41_series_optimized': True,
                'gpt41_large_context_ready': True,
                'supports_1m_tokens': True
            }
            
            if filename not in self.document_metadata.get('processed_files', []):
                self.document_metadata['processed_files'].append(filename)
            
            self.document_metadata['file_info'][filename] = file_metadata
            self.document_metadata['last_update'] = datetime.now().isoformat()
            self.document_metadata['total_chunks'] = self.document_metadata.get('total_chunks', 0) + len(chunks)
            
            metadata_file = os.path.join(self.vector_store_path, "document_metadata.json")
            with open(metadata_file, 'w') as f:
                json.dump(self.document_metadata, f, indent=2)
            
            return ProcessingResult(
                success=True,
                filename=filename,
                chunks_added=len(chunks),
                document_type=file_ext,
                processing_time=processing_time,
                extraction_methods=[processor_func.__name__],
                metadata=file_metadata,
                academic_content_detected=academic_content_detected
            )
            
        except Exception as e:
            processing_time = time.time() - start_time
            
            return ProcessingResult(
                success=False,
                filename=filename,
                chunks_added=0,
                document_type=os.path.splitext(filename)[1][1:] if '.' in filename else 'unknown',
                processing_time=processing_time,
                extraction_methods=["failed"],
                metadata={},
                error=str(e),
                academic_content_detected=False
            )
    
    def _update_system_metadata(self, processing_results: Dict[str, Any]):
        """Update system-wide metadata after processing"""
        try:
            for subject in processing_results.get('subjects_identified', []):
                if subject not in self.document_metadata.get('subjects', {}):
                    self.document_metadata.setdefault('subjects', {})[subject] = []
                
                self.document_metadata['subjects'][subject].extend(
                    processing_results['categories_detected'].get(subject, [])
                )
            
            for filename in processing_results.get('successful_files', []):
                if filename in self.document_metadata.get('file_info', {}):
                    file_info = self.document_metadata['file_info'][filename]
                    category_info = file_info.get('category', {})
                    academic_level = category_info.get('academic_level')
                    
                    if academic_level:
                        self.document_metadata.setdefault('academic_levels', {})[academic_level] = \
                            self.document_metadata.get('academic_levels', {}).get(academic_level, 0) + 1
            
            # NEW: Add GPT-4.1 system metadata
            self.document_metadata['gpt41_series'] = True
            self.document_metadata['supports_1m_tokens'] = True
            self.document_metadata['last_gpt41_update'] = datetime.now().isoformat()
            
            metadata_file = os.path.join(self.vector_store_path, "document_metadata.json")
            with open(metadata_file, 'w') as f:
                json.dump(self.document_metadata, f, indent=2)
                
        except Exception as e:
            pass

    def _build_enhanced_search_indices(self):
        """GPT-4.1 ENHANCED: Build search indices"""
        try:
            self._build_tfidf_index()
        except Exception as e:
            pass

    def _build_tfidf_index(self):
        """GPT-4.1 ENHANCED: Build TF-IDF index"""
        try:
            if not self.vectorstore or self.vectorstore.index.ntotal <= 1:
                return
            
            with self.tfidf_lock:
                try:
                    all_docs = self.vectorstore.similarity_search("", k=15000)
                    self.tfidf_chunks = [doc for doc in all_docs 
                                       if doc.metadata.get('source') != 'init']
                except Exception:
                    return
                
                if not self.tfidf_chunks:
                    return
                
                texts = []
                for doc in self.tfidf_chunks:
                    enhanced_text = doc.page_content
                    
                    metadata = doc.metadata
                    if metadata.get('primary_subject'):
                        enhanced_text += f" {metadata['primary_subject']}"
                    if metadata.get('chunk_academic_type'):
                        enhanced_text += f" {metadata['chunk_academic_type']}"
                    if metadata.get('gpt41_optimized'):
                        enhanced_text += " gpt41_optimized"
                    
                    texts.append(enhanced_text.lower())
                
                self.tfidf_matrix = self.tfidf_vectorizer.fit_transform(texts)
                
        except Exception as e:
            pass
        
    def rebuild_vectorstore(self) -> Dict[str, Any]:
        """GPT-4.1 ENHANCED: Vectorstore rebuild with robust error handling"""
        start_time = time.time()
        
        try:
            def safe_remove_directory(path, max_retries=5, wait_time=1.0):
                """Safely remove directory with retry logic"""
                for attempt in range(max_retries):
                    try:
                        if os.path.exists(path):
                            if hasattr(self, 'vectorstore') and self.vectorstore:
                                try:
                                    del self.vectorstore
                                    self.vectorstore = None
                                except Exception:
                                    pass
                            
                            time.sleep(0.5)
                            
                            for root, dirs, files in os.walk(path, topdown=False):
                                for file in files:
                                    file_path = os.path.join(root, file)
                                    try:
                                        os.chmod(file_path, 0o777)
                                        os.remove(file_path)
                                    except Exception as e:
                                        pass
                                
                                for dir_name in dirs:
                                    dir_path = os.path.join(root, dir_name)
                                    try:
                                        os.rmdir(dir_path)
                                    except Exception as e:
                                        pass
                            
                            os.rmdir(path)
                            return True
                            
                    except Exception as e:
                        if "Device or resource busy" in str(e) or "resource busy" in str(e):
                            time.sleep(wait_time)
                            wait_time *= 1.5
                        else:
                            if attempt == max_retries - 1:
                                return False
                
                return False
            
            with self.processing_lock:
                if hasattr(self, 'vectorstore') and self.vectorstore:
                    try:
                        del self.vectorstore
                        self.vectorstore = None
                    except Exception as e:
                        pass
                
                removal_success = safe_remove_directory(self.vector_store_path, max_retries=5)
                
                if not removal_success:
                    try:
                        backup_name = f"{self.vector_store_path}_backup_{int(time.time())}"
                        os.rename(self.vector_store_path, backup_name)
                    except Exception as e:
                        pass
                
                os.makedirs(self.vector_store_path, exist_ok=True)
                
                self.document_metadata = {
                    'processed_files': [],
                    'last_update': datetime.now().isoformat(),
                    'total_chunks': 0,
                    'file_info': {},
                    'categories': {},
                    'subjects': {},
                    'academic_levels': {},
                    'gpt41_series': True,
                    'supports_1m_tokens': True,
                    'optimized_response_lengths': True
                }
                
                self.vectorstore = None

            uploads_dir = 'uploads'
            errors = []
            
            if not os.path.exists(uploads_dir):
                return {
                    'success': True,
                    'processed_files': [],
                    'failed_files': [],
                    'total_chunks': 0,
                    'processing_time': time.time() - start_time,
                    'errors': ['No uploads directory found'],
                    'message': 'No uploads directory found',
                    'gpt41_series': True
                }

            file_paths = []
            for filename in os.listdir(uploads_dir):
                filepath = os.path.join(uploads_dir, filename)
                if os.path.isfile(filepath):
                    file_ext = os.path.splitext(filename)[1].lower()
                    if file_ext in self.document_processors:
                        file_paths.append(filepath)
                    else:
                        errors.append(f"Unsupported file type: {filename}")

            if not file_paths:
                return {
                    'success': True,
                    'processed_files': [],
                    'failed_files': [],
                    'total_chunks': 0,
                    'processing_time': time.time() - start_time,
                    'errors': ['No supported files found'],
                    'message': 'No supported files found',
                    'gpt41_series': True
                }

            try:
                processing_results = self.process_documents(file_paths)
                processing_time = time.time() - start_time
                
                for failed_file in processing_results.get('failed_files', []):
                    if isinstance(failed_file, dict) and 'error' in failed_file:
                        errors.append(f"{failed_file.get('filename', 'unknown')}: {failed_file['error']}")

                return {
                    'success': True,
                    'processed_files': processing_results.get('successful_files', []),
                    'failed_files': processing_results.get('failed_files', []),
                    'total_chunks': processing_results.get('total_chunks_added', 0),
                    'processing_time': processing_time,
                    'subjects_identified': processing_results.get('subjects_identified', []),
                    'categories_detected': processing_results.get('categories_detected', {}),
                    'academic_documents': processing_results.get('academic_content_detected', 0),
                    'enhanced_features': True,
                    'large_document_support': True,
                    'optimized_response_lengths': True,
                    'gpt41_series': True,
                    'supports_1m_tokens': True,
                    'gpt41_models': ['gpt-4.1-nano', 'gpt-4.1-mini', 'gpt-4.1'],
                    'errors': errors
                }

            except Exception as processing_error:
                processing_time = time.time() - start_time
                error_msg = f"GPT-4.1 Processing failed: {str(processing_error)}"
                errors.append(error_msg)
                
                return {
                    'success': False,
                    'processed_files': [],
                    'failed_files': [],
                    'total_chunks': 0,
                    'processing_time': processing_time,
                    'errors': errors,
                    'message': 'GPT-4.1 Processing failed',
                    'gpt41_series': True
                }

        except Exception as e:
            processing_time = time.time() - start_time
            error_msg = f"GPT-4.1 Vectorstore rebuild failed: {str(e)}"
            
            return {
                'success': False,
                'processed_files': [],
                'failed_files': [],
                'total_chunks': 0,
                'processing_time': processing_time,
                'errors': [error_msg],
                'message': 'GPT-4.1 Rebuild failed',
                'gpt41_series': True
            }
        
    def _initialize_vectorstore(self):
        """GPT-4.1 ENHANCED: Initialize vectorstore"""
        try:
            vectorstore_path = os.path.join(self.vector_store_path, "index.faiss")
            if os.path.exists(vectorstore_path):
                try:
                    self.vectorstore = FAISS.load_local(
                        self.vector_store_path,
                        self.embeddings,
                        allow_dangerous_deserialization=True
                    )

                    metadata_file = os.path.join(self.vector_store_path, "document_metadata.json")
                    if os.path.exists(metadata_file):
                        with open(metadata_file, 'r') as f:
                            self.document_metadata.update(json.load(f))

                    self._build_enhanced_search_indices()
                    return
                except Exception as e:
                    # If loading fails, set vectorstore to None and continue
                    self.vectorstore = None
            else:
                # No existing vectorstore, initialize as None
                self.vectorstore = None
                
        except Exception as e:
            # If any error occurs, ensure vectorstore is None
            self.vectorstore = None

    # ALL COMPATIBILITY METHODS
    
    def is_file_in_vectorstore(self, filename: str) -> bool:
        """Check if file is already processed"""
        return filename in self.document_metadata.get('processed_files', [])

    def clear_conversation_memory(self, session_id: str):
        """Clear conversation memory for session"""
        with self.conversation_lock:
            if session_id in self.conversation_contexts:
                del self.conversation_contexts[session_id]

    def get_processed_files(self) -> List[Dict[str, Any]]:
        """GPT-4.1 ENHANCED: Get processed files with enhanced metadata"""
        processed_files = []
        
        try:
            file_info = self.document_metadata.get('file_info', {})
            for filename, info in file_info.items():
                category_info = info.get('category', {})
                processed_files.append({
                    'filename': filename,
                    'document_type': info.get('file_type', 'unknown'),
                    'chunks': info.get('chunk_count', 0),
                    'academic_content': info.get('academic_content_detected', False),
                    'processing_timestamp': info.get('processed_at', ''),
                    'primary_subject': category_info.get('primary_subject', 'unknown'),
                    'academic_level': category_info.get('academic_level', 'unknown'),
                    'confidence_score': category_info.get('confidence_score', 0.0),
                    'key_topics': category_info.get('key_topics', []),
                    'large_document': info.get('large_document_processing', False),
                    'optimized_response_system': info.get('optimized_response_system', True),
                    'gpt41_series_optimized': info.get('gpt41_series_optimized', True),
                    'supports_1m_tokens': info.get('supports_1m_tokens', True)
                })
        except Exception as e:
            pass
        
        return processed_files

    def get_routing_stats(self) -> Dict[str, Any]:
        """GPT-4.1 ENHANCED: Get comprehensive routing and system statistics"""
        return {
            "routing_performance": self.routing_stats,
            "document_categories": self.document_metadata.get('subjects', {}),
            "academic_levels": self.document_metadata.get('academic_levels', {}),
            "total_documents": len(self.document_metadata.get('processed_files', [])),
            "conversation_sessions": len(self.conversation_contexts),
            "enhanced_features_active": True,
            "large_document_support": True,
            "optimized_response_lengths": True,
            "gpt41_series": True,
            "supports_1m_tokens": True,
            "last_update": self.document_metadata.get('last_update', ''),
            "gpt41_model_usage": {
                "nano_queries": self.routing_stats.get("gpt41_nano_queries", 0),
                "mini_queries": self.routing_stats.get("gpt41_mini_queries", 0),
                "full_queries": self.routing_stats.get("gpt41_full_queries", 0)
            },
            "system_capabilities": {
                "multi_level_routing": True,
                "conversation_memory": True,
                "advanced_categorization": True,
                "intelligent_context_scoring": True,
                "dynamic_document_handling": True,
                "large_document_processing": True,
                "adaptive_response_length": True,
                "explicit_detail_detection": True,
                "concise_simple_responses": True,
                "comprehensive_detail_responses": True,
                "gpt41_nano_ultra_fast": True,
                "gpt41_mini_balanced": True,
                "gpt41_full_comprehensive": True,
                "massive_context_1m_tokens": True
            }
        }

    def get_system_stats(self) -> Dict[str, Any]:
        """GPT-4.1 ENHANCED: Get comprehensive system statistics"""
        return {
            "vectorstore": {
                "total_vectors": self.vectorstore.index.ntotal if self.vectorstore else 0,
                "total_documents": len(self.document_metadata.get('processed_files', [])),
                "subjects": list(self.document_metadata.get('subjects', {}).keys()),
                "academic_levels": list(self.document_metadata.get('academic_levels', {}).keys())
            },
            "enhanced_features": {
                "conversation_memory": True,
                "multi_level_routing": True,
                "advanced_categorization": True,
                "intelligent_relevance_scoring": True,
                "dynamic_chunking": True,
                "academic_type_detection": True,
                "large_document_processing": True,
                "adaptive_response_length": True,
                "optimized_length_control": True,
                "explicit_detail_request_detection": True
            },
            "gpt41_capabilities": {
                "series": "GPT-4.1",
                "models_available": ["gpt-4.1-nano", "gpt-4.1-mini", "gpt-4.1"],
                "max_context_tokens": 1000000,
                "ultra_fast_nano": True,
                "balanced_mini": True,
                "comprehensive_full": True,
                "optimized_chunking": True,
                "massive_context_support": True
            },
            "performance": {
                "parallel_processing": self.enable_parallel,
                "max_workers": self.max_workers,
                "conversation_sessions": len(self.conversation_contexts),
                "embedding_batch_size": self.embedding_batch_size,
                "max_tokens_per_batch": self.max_tokens_per_batch,
                "large_doc_threshold": self.large_doc_threshold
            },
            "optimization_stats": {
                "concise_responses_enabled": True,
                "detailed_explanation_detection": True,
                "explicit_detail_requests": self.routing_stats.get("explicit_detail_requests", 0),
                "response_length_optimization": True,
                "gpt41_nano_usage": self.routing_stats.get("gpt41_nano_queries", 0),
                "gpt41_mini_usage": self.routing_stats.get("gpt41_mini_queries", 0),
                "gpt41_full_usage": self.routing_stats.get("gpt41_full_queries", 0)
            }
        }

# BACKWARD COMPATIBILITY
DocumentProcessor = UltraRobustAcademicDocumentProcessor

# EXPORT
__all__ = ['UltraRobustAcademicDocumentProcessor', 'DocumentProcessor']

if __name__ == "__main__":
    processor = UltraRobustAcademicDocumentProcessor()
